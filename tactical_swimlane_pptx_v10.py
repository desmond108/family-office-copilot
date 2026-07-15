"""tactical_swimlane_pptx.py — the v10 tactical-instructions process as a native
16:9 PowerPoint swimlane, matching tactical_swimlane_note.html (v10).

Two slides:
  1. The answer — how do the client's tactical instructions reach the proposal?
     (v10: they pass VERBATIM into one self-contained prompt; the AI Model writes
     the narrative; every number stays deterministic).
  2. The swimlane — three actor lanes (Client / You / Copilot) across four phases:
     Set policy -> Capture -> Compute (deterministic) -> Assemble, generate & deliver.

House navy / gold, matching the other decks. Run:
    python tactical_swimlane_pptx.py   ->  Tactical_Swimlane_v10.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1E, 0x2A, 0x56); NAVY_DEEP = RGBColor(0x13, 0x1D, 0x3D)
GOLD = RGBColor(0xB0, 0x87, 0x2A); GOLD_DK = RGBColor(0x9C, 0x74, 0x22)
GOLD_TINT = RGBColor(0xF7, 0xEF, 0xDC)
GREEN = RGBColor(0x2F, 0x6F, 0x57); GREEN_TINT = RGBColor(0xE6, 0xF1, 0xEB)
AMBER = RGBColor(0xB0, 0x6A, 0x17); AMBER_TINT = RGBColor(0xFB, 0xE6, 0xD2)
INK = RGBColor(0x16, 0x20, 0x3C); SOFT = RGBColor(0x47, 0x51, 0x7A)
FAINT = RGBColor(0x83, 0x8C, 0xAD); LINE = RGBColor(0xDF, 0xE2, 0xEE)
LINE_STRONG = RGBColor(0xC8, 0xCD, 0xDD); SLATE_TINT = RGBColor(0xEE, 0xF0, 0xF6)
CLIENT_TINT = RGBColor(0xEE, 0xF2, 0xFB)
PAPER = RGBColor(0xFF, 0xFF, 0xFF); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xB9, 0xC2, 0xE0)

SERIF = "Georgia"; SANS = "Calibri"; MONO = "Consolas"

prs = Presentation()
prs.slide_width = In(13.333); prs.slide_height = In(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08, dash=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        if dash:  # solidFill is already set above, so prstDash is in schema order
            from pptx.oxml.ns import qn
            ln = shp.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = In(0.08); tf.margin_right = In(0.07)
    tf.margin_top = In(0.05); tf.margin_bottom = In(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    try:
        tf.auto_size = None
    except Exception:
        pass
    return shp


def txt(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    try:
        tf.auto_size = None
    except Exception:
        pass
    return tf


def para(tf, runs, *, first=False, align=PP_ALIGN.LEFT, before=0, after=0, spacing=1.0):
    """runs: list of (text, size, color, bold, font). Single paragraph, multi-run."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(before); p.space_after = Pt(after)
    try:
        p.line_spacing = spacing
    except Exception:
        pass
    for text, size, color, bold, font in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font; r.font.color.rgb = color
    return p


def one(tf, text, size, color, **kw):
    return para(tf, [(text, size, color, kw.pop("bold", False), kw.pop("font", SANS))], **kw)


# =========================================================================== #
# Slide 1 — how tactical instructions reach the proposal (v10)
# =========================================================================== #
s = slide()
one(txt(s, 0.8, 0.5, 11.5, 0.4), "MERIDIAN FAMILY OFFICE COPILOT · v10", 11.5, GOLD,
    bold=True, first=True)
one(txt(s, 0.8, 0.92, 11.9, 0.9),
    "How do the client’s tactical instructions reach the proposal?", 26, NAVY, font=SERIF,
    bold=True, first=True)

cards = [
    ("Pass through — verbatim", "no sorting, no classifying", NAVY,
     "The client’s tactical text goes straight into the prompt, unedited — together with the "
     "intake parameters, the parsed holdings + statement source, and the research / other "
     "documents in full.", SLATE_TINT, LINE_STRONG, INK),
    ("One self-contained prompt", "the transparency surface", GOLD_DK,
     "Shown on the Proposal page — editable, copyable, downloadable. Paste it into ANY AI model "
     "to reproduce the proposal and compare outputs. Generate in-app runs the live AI model.",
     GOLD_TINT, GOLD, GOLD_DK),
    ("Numbers stay deterministic", "the guardrail", GREEN,
     "Every figure is computed by the engine (the FACTS block). The AI writes prose only; "
     "qualitative claims from the documents are context, not independently verified.",
     GREEN_TINT, GREEN, INK),
]
cy = 2.25; ch = 1.25; gap = 0.24
for i, (title, kicker, tcol, body, fill, lcol, bodycol) in enumerate(cards):
    y = cy + i * (ch + gap)
    b = rect(s, 0.8, y, 11.73, ch, fill=fill, line=lcol, line_w=1.5, radius=0.06)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(f"{i+1}.  ", 15, tcol, True, SANS), (title, 15, tcol, True, SANS),
              (f"   — {kicker}", 12, SOFT, False, SANS)], first=True)
    one(tf, body, 12.5, bodycol, before=5, spacing=1.08)

b = rect(s, 0.8, 6.5, 11.73, 0.95, fill=NAVY_DEEP, radius=0.06)
tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, [("What changed in v10:  ", 11.5, GOLD, True, SANS),
          ("the copilot no longer sorts instructions into typed items, tags enforcement tiers, or "
           "gates a trade on a price trigger. It hands the raw inputs to the AI Model in one prompt "
           "and lets it analyse them — while the engine still computes every number, so nothing is "
           "invented. The UI is model-agnostic: it always reads “AI Model”.", 11.5, WHITE, False, SANS)],
     first=True, spacing=1.04)


# =========================================================================== #
# Slide 2 — the swimlane (v10)
# =========================================================================== #
s = slide()
one(txt(s, 0.5, 0.34, 11.5, 0.4), "MERIDIAN FAMILY OFFICE COPILOT · v10 · PROCESS",
    11, GOLD, bold=True, first=True)
one(txt(s, 0.5, 0.68, 12.4, 0.6), "Tactical instructions — who does what, and when",
    23, NAVY, font=SERIF, bold=True, first=True)

# geometry
LX, LW = 0.4, 0.95
GX = LX + LW                      # 1.35
GR = 13.0
PH_Y, PH_H = 1.5, 0.38
# phase widths (D widest — most action now lives in assemble/generate/deliver)
PW = [2.7, 2.4, 2.75, 3.8]
PXS = [GX]
for w in PW:
    PXS.append(PXS[-1] + w)      # PXS[i]..PXS[i+1] is phase i

LANES = [("Client", "Source", CLIENT_TINT, 1.92, 1.18),
         ("You", "Analyst", RGBColor(0xFB, 0xF6, 0xE9), 3.16, 2.46),
         ("Copilot", "System", SLATE_TINT, 5.68, 1.5)]

# phase headers
phase_titles = [("A", "Set profile & policy"), ("B", "Capture inputs"),
                ("C", "Compute · deterministic"), ("D", "Assemble → generate → deliver")]
rect(s, LX, PH_Y, LW, PH_H, fill=NAVY_DEEP, radius=0.04)
for i, (pn, pt) in enumerate(phase_titles):
    b = rect(s, PXS[i], PH_Y, PW[i], PH_H, fill=NAVY_DEEP, radius=0.04)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(f"{pn}  ", 11, GOLD, True, MONO), (pt, 11, WHITE, True, SANS)],
         first=True, align=PP_ALIGN.CENTER)

# lane bands + labels
for name, sub, bg, y, h in LANES:
    rect(s, GX, y, GR - GX, h, fill=bg, line=LINE, line_w=0.75, radius=0.01)
    lb = rect(s, LX, y, LW, h, fill=None, line=None)
    tf = lb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    col = {"You": GOLD_DK, "Client": NAVY, "Copilot": SOFT}[name]
    para(tf, [(name, 13, col, True, SANS)], first=True, align=PP_ALIGN.CENTER)
    para(tf, [(sub.upper(), 8.5, FAINT, False, SANS)], align=PP_ALIGN.CENTER, before=1)

CY = {name: (y, h) for name, _, _, y, h in LANES}


def abox(px, y, w, h, kind, num, title, desc, tsize=10, dsize=8.3):
    fill, lcol, lw, dash, ncol = {
        "you": (PAPER, GOLD, 1.25, False, GOLD_DK),
        "bot": (PAPER, LINE_STRONG, 1.0, True, FAINT),
        "cli": (PAPER, RGBColor(0x9C, 0xA8, 0xCE), 1.0, False, NAVY),
    }[kind]
    b = rect(s, px, y, w, h, fill=fill, line=lcol, line_w=lw, radius=0.10, dash=dash)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    head = []
    if num:
        head.append((f"{num} · ", tsize, ncol, True, SANS))
    head.append((title, tsize, INK if kind != "bot" else SOFT, True, SANS))
    para(tf, head, first=True, spacing=1.0)
    if desc:
        one(tf, desc, dsize, SOFT, before=2, spacing=1.02)
    return b


def gear(px, y, w, h, tag, title, desc, accent=False):
    lcol = GOLD if accent else LINE_STRONG
    b = rect(s, px, y, w, h, fill=PAPER, line=lcol, line_w=1.3 if accent else 1.0,
             radius=0.10, dash=not accent)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    one(tf, tag, 8, GOLD_DK if accent else FAINT, bold=True, first=True)
    one(tf, title, 10, INK if accent else SOFT, bold=True, before=2, spacing=1.0)
    if desc:
        one(tf, desc, 8.3, SOFT, before=2, spacing=1.0)
    return b

# ---- CLIENT lane ---- #
cy, ch = CY["Client"]
# B: gives instructions + documents
abox(PXS[1] + 0.12, cy + 0.2, PW[1] - 0.24, ch - 0.4, "cli", None,
     "Gives instructions + documents", "“gold below $4,000” · statements · research",
     tsize=10.5, dsize=8)

# ---- YOU lane ---- #
yy, yh = CY["You"]
# Phase A: boxes 1 (top) + 2 (bottom) — set the policy first
ax, aw = PXS[0] + 0.12, PW[0] - 0.24
abox(ax, yy + 0.24, aw, 1.02, "you", "1", "Set mandate, risk & ability", None)
abox(ax, yy + 1.34, aw, 1.02, "you", "2", "Set allocation targets & limits", "limits always manual")
# Phase B: box 3 — paste + upload (verbatim; no sorting)
abox(PXS[1] + 0.12, yy + 0.6, PW[1] - 0.24, 1.3, "you", "3",
     "Paste tactical text (verbatim) + upload documents", "the client’s words, unedited",
     tsize=10, dsize=8)
# Phase C: box 4 — Analyse, after the policy is set
abox(PXS[2] + 0.12, yy + 0.75, PW[2] - 0.24, 0.95, "you", "4",
     "Analyse ▸", "after the policy above is set")
# Phase D: boxes 5, 6, 7 stacked
dx, dw = PXS[3] + 0.12, PW[3] - 0.24
abox(dx, yy + 0.16, dw, 0.66, "you", "5", "Review the assembled prompt", None)
abox(dx, yy + 0.90, dw, 0.66, "you", "6", "Generate — or copy to any AI model", None)
abox(dx, yy + 1.64, dw, 0.66, "you", "7", "Review & download the deck", "PPTX / PDF")

# ---- COPILOT lane ---- #
py, ph = CY["Copilot"]
# C: parse + compute (stacked, narrow lane)
gear(PXS[2] + 0.12, py + 0.14, PW[2] - 0.24, 0.58, "⚙ ON ANALYSE",
     "Parse → build the book", None)
gear(PXS[2] + 0.12, py + 0.78, PW[2] - 0.24, 0.58, "⚙ AUTO",
     "Compute deterministic FACTS", "allocation · drift · suitability")
# D: assemble → AI model → render (three across)
gw = (PW[3] - 0.24 - 2 * 0.12) / 3
g1x = PXS[3] + 0.12
g2x = g1x + gw + 0.12
g3x = g2x + gw + 0.12
gear(g1x, py + 0.3, gw, 0.9, "⚙ ASSEMBLE", "One self-contained prompt", None)
gear(g2x, py + 0.3, gw, 0.9, "✨ AI MODEL", "Writes the narrative", None, accent=True)
gear(g3x, py + 0.3, gw, 0.9, "⚙ RENDER", "PPTX / PDF deck", "figures deterministic")

# footer key point
one(txt(s, 0.5, 7.16, 12.5, 0.3),
    "Recommended order: set the policy (1–2) first, then Analyse (4). The tactical text and documents "
    "pass verbatim into one self-contained prompt (5–6) that the AI Model writes from — and that you can "
    "copy into any AI model. Every figure is computed deterministically; the AI writes prose only.",
    10, SOFT, bold=True, first=True)

OUT = "Tactical_Swimlane_v10.pptx"
prs.save(OUT)
print(f"wrote {OUT} — {len(prs.slides._sldIdLst)} slides")
