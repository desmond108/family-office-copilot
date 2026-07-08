"""generate_proposal.py — build a proposal deck from the PARSED client book.

The Proposal page no longer shows a static example: it renders a deck generated
from whatever statements the analyst loaded (e.g. the three sample custodian
statements), in the house navy / gold format, and offers the SAME deck as a
PowerPoint or PDF download.

One computed model (built in app.py from the deterministic engine) drives three
renderers so the on-screen deck and both downloads always agree:

    render_html(model)  -> str    (inline, navy/gold, reuses proposal_deck CSS)
    render_pptx(model)  -> bytes  (python-pptx)
    render_pdf(model)   -> bytes  (reportlab)

Every figure in ``model`` is computed by suitability_check / the parser / the
rebalancer — this module only lays them out. Nothing here invents a number.
"""
from __future__ import annotations

import html
from io import BytesIO

from proposal_deck import STYLE, _slide  # reuse the house navy/gold CSS

# --------------------------------------------------------------------------- #
# Shared palette (kept in sync with proposal_deck.py)
# --------------------------------------------------------------------------- #
NAVY = "#1e2a56"
GOLD = "#b0872a"
SOFT = "#5a648a"
RULE = "#dfe4ef"
OK = "#2f7a52"
WARN = "#b3402f"
ENFORCE = {"block": WARN, "flag": "#b07a1e", "disclose": OK, "none": OK}


# --------------------------------------------------------------------------- #
# 1) HTML renderer (inline view)
# --------------------------------------------------------------------------- #
def _html_table(headers, rows, aligns=None, total_row=None):
    aligns = aligns or ["left"] * len(headers)
    th = "".join(f'<th style="text-align:{a}">{html.escape(h)}</th>'
                 for h, a in zip(headers, aligns))
    body = ""
    for r in rows:
        body += "<tr>" + "".join(
            f'<td style="text-align:{a}">{c}</td>' for c, a in zip(r, aligns)) + "</tr>"
    if total_row:
        body += '<tr class="tot">' + "".join(
            f'<td style="text-align:{a}">{c}</td>' for c, a in zip(total_row, aligns)) + "</tr>"
    return f'<table class="tbl"><tr>{th}</tr>{body}</table>'


def _stat(lab, big):
    return (f'<div class="stat"><div class="lab">{html.escape(lab)}</div>'
            f'<div class="big">{html.escape(big)}</div></div>')


def _narr_paras(text: str) -> str:
    """Split a narrative string into <p> paragraphs for the HTML deck."""
    return "".join(f'<p class="note" style="font-size:14px;line-height:1.6">'
                   f'{html.escape(p.strip())}</p>'
                   for p in text.split("\n\n") if p.strip())


def render_html(m: dict) -> str:
    meta = m["meta"]
    ent = "entity" if meta["entities"] == 1 else "entities"
    off = 1 if m.get("narrative") else 0

    s1 = _slide(1, "", "", (
        '<div class="eyebrow">Portfolio Proposal · Confidential</div>'
        f'<h1>{html.escape(m["title"])}</h1>'
        f'<div class="sub">{html.escape(m["subtitle"])}</div>'
        f'<div class="lede">Generated from {meta["custodians"]} parsed custodian statement(s) '
        f'across {meta["entities"]} {ent} · {meta["positions"]} positions. Every figure is '
        f'computed by the deterministic engine from the client\'s own holdings — nothing is '
        f'invented.</div>'
        f'<div class="stamp">As of {html.escape(str(m["as_of"]))}</div>'), cover=True)

    snarr = ""
    if m.get("narrative"):
        snarr = _slide(2, "Investment Commentary", "Chief Investment Office — Commentary", (
            _narr_paras(m["narrative"])
            + '<p class="note" style="color:#9AA4C4"><b>Note:</b> This commentary is '
            'generated prose grounded strictly in the deterministic figures shown in the '
            'following slides; it quotes those figures but does not compute or alter any of '
            'them. For discussion only; not investment advice.</p>'))

    stats = "".join(_stat(l, v) for l, v in m["metrics"])
    s2 = _slide(2 + off, "Position", "Current Consolidated Position", (
        f'<div class="grid2">{stats}</div>'
        f'<p class="note"><b>Mandate</b> {html.escape(meta["mandate"])} · <b>Risk appetite</b> '
        f'{html.escape(meta["risk"])} · <b>Ability to take risk</b> {html.escape(meta["ability"])}. '
        f'Consolidated across {html.escape(", ".join(m["custodian_list"]))}. '
        f'{html.escape(m["prov"])}</p>'))

    atable = _html_table(
        ["Asset class", "Value (USD)", "% of gross", "Target", "Drift"],
        m["alloc_rows"], aligns=["left", "right", "right", "right", "right"],
        total_row=["Gross assets", m["gross_str"], "100.0%", "", ""])
    s3 = _slide(3 + off, "Allocation", "Current Allocation vs Target", (
        atable
        + '<p class="note">Weights are computed from parsed statement values against the '
        'mandate target. Sleeves with <b>no target</b> (e.g. alternatives, real estate) are '
        'flagged on the suitability slide, not treated as compliant by omission.</p>'))

    rtable = _html_table(
        ["Sleeve", "Before", "Target", "Trade (USD)", "Action", "Note"],
        m["reb_rows"], aligns=["left", "right", "right", "right", "left", "left"])
    rs = m["reb_summary"]
    fund = ("nets to ≈$0 — self-funding" if rs["selffund"]
            else "requires external cash / raises proceeds")
    s4 = _slide(4 + off, "Proposal", "Rebalancing Proposal — Before → After", (
        rtable
        + f'<p class="note"><b>Buys</b> {rs["buys"]} · <b>Sells</b> {rs["sells"]} · '
        f'<b>Net of trades</b> {rs["net"]} ({fund}). Illiquid sleeves cannot be traded on '
        f'demand — either amend the mandate to add them to target, or stage the reduction '
        f'against redemption windows.</p>'))

    if m["suit_items"]:
        rows = "".join(
            f'<div style="display:flex;gap:10px;padding:9px 0;border-bottom:1px solid {RULE}">'
            f'<span style="font-family:monospace;font-size:11px;font-weight:700;color:#fff;'
            f'background:{ENFORCE.get(e, SOFT)};padding:2px 9px;border-radius:5px;'
            f'text-transform:uppercase;height:fit-content">{html.escape(e)}</span>'
            f'<span style="font-size:13.5px;color:{NAVY}">{html.escape(d)}</span></div>'
            for e, d in m["suit_items"])
    else:
        rows = '<p class="note">Book is within all defined suitability bands.</p>'
    s5 = _slide(5 + off, "Suitability", "Suitability of the Proposed Book", (
        f'<p class="note" style="margin-top:0">Mandate gate · '
        f'<b>{html.escape(m["gate"])}</b></p>{rows}'))

    if m["data_quality"]:
        dq = "".join(
            f'<div style="display:flex;gap:10px;padding:9px 0;border-bottom:1px solid {RULE}">'
            f'<span style="font-family:monospace;font-size:11px;font-weight:700;color:#fff;'
            f'background:#b07a1e;padding:2px 9px;border-radius:5px">{html.escape(k)}</span>'
            f'<span style="font-size:13.5px;color:{NAVY}">{html.escape(d)}</span></div>'
            for k, d in m["data_quality"])
    else:
        dq = '<p class="note">No data-quality issues detected.</p>'

    notes_html = ""
    if m["overlays"]:
        notes_html += ('<p class="note" style="margin-top:0"><b>Overlay sleeves (analyst):</b> '
                       + html.escape(" · ".join(m["overlays"])) + "</p>")
    if m["analyst_notes"]:
        notes_html += '<ul class="reqs" style="counter-reset:none;list-style:disc;padding-left:20px">'
        notes_html += "".join(f'<li style="padding-left:0">{html.escape(n)}</li>'
                              for n in m["analyst_notes"])
        notes_html += "</ul>"
    if not notes_html:
        notes_html = '<p class="note" style="margin-top:0">No analyst notes captured.</p>'

    s6 = _slide(6 + off, "Data & Method", "Data Quality, Analyst Notes & Provenance", (
        '<div class="two"><div><h3>Reconciliation & data-quality flags</h3>' + dq + "</div>"
        '<div><h3>Analyst notes folded into this proposal</h3>' + notes_html
        + f'<p class="note">{html.escape(m["prov"])} Figures are computed by the deterministic '
        f'suitability / rebalancing engine; unresolved ids and stale marks are flagged, never '
        f'smoothed. For discussion only; not investment advice.</p></div></div>'))

    slides = [s1] + ([snarr] if snarr else []) + [s2, s3, s4, s5, s6]
    return STYLE + '<div class="wrap">' + "".join(slides) + "</div>"


# --------------------------------------------------------------------------- #
# 2) PPTX renderer
# --------------------------------------------------------------------------- #
def render_pptx(m: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    def rgb(hexs):
        return RGBColor.from_string(hexs.lstrip("#").upper())

    NAVY_C, GOLD_C, SOFT_C, WHITE_C = rgb(NAVY), rgb(GOLD), rgb(SOFT), rgb("#FFFFFF")
    off = 1 if m.get("narrative") else 0
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tb.text_frame.word_wrap = True
        tb.text_frame.vertical_anchor = anchor
        return tb.text_frame

    def para(tf, text, size, color, bold=False, first=False, space=4):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_after = Pt(space)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        return p

    def header(slide, eyebrow, title):
        tf = textbox(slide, Inches(0.6), Inches(0.42), Inches(12), Inches(0.4))
        para(tf, eyebrow.upper(), 12, GOLD_C, bold=True, first=True)
        tf2 = textbox(slide, Inches(0.6), Inches(0.82), Inches(12.1), Inches(0.9))
        para(tf2, title, 30, NAVY_C, bold=True, first=True)

    def footer(slide, n):
        tf = textbox(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35))
        p = para(tf, "Portfolio Proposal · Confidential · computed by the deterministic engine",
                 9, rgb("#9AA4C4"), first=True)
        p.alignment = PP_ALIGN.LEFT

    def table(slide, headers, rows, left, top, width, col_ratios=None, total=None,
              aligns=None, fs=11):
        nrows = len(rows) + 1 + (1 if total else 0)
        gt = slide.shapes.add_table(nrows, len(headers), left, top, width,
                                    Inches(0.34 * nrows)).table
        gt.first_row = False
        gt.horz_banding = False
        if col_ratios:
            tot = sum(col_ratios)
            for i, cr in enumerate(col_ratios):
                gt.columns[i].width = int(width * cr / tot)
        aligns = aligns or [PP_ALIGN.LEFT] * len(headers)
        allrows = [headers] + rows + ([total] if total else [])
        for ri, row in enumerate(allrows):
            is_head = ri == 0
            is_total = total is not None and ri == len(allrows) - 1
            for ci, val in enumerate(row):
                cell = gt.cell(ri, ci)
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.margin_top = Inches(0.03)
                cell.margin_bottom = Inches(0.03)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                if is_head:
                    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY_C
                elif is_total:
                    cell.fill.solid(); cell.fill.fore_color.rgb = rgb("#F4F6FB")
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = WHITE_C
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = aligns[ci]
                run = p.add_run(); run.text = str(val)
                run.font.size = Pt(fs)
                run.font.bold = is_head or is_total
                run.font.name = "Calibri"
                run.font.color.rgb = WHITE_C if is_head else NAVY_C
        return gt

    # ---- Slide 1 · cover -------------------------------------------------- #
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY_C
    tf = textbox(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
    para(tf, "PORTFOLIO PROPOSAL · CONFIDENTIAL", 13, rgb("#C8A24A"), bold=True, first=True)
    tf = textbox(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.5))
    para(tf, m["title"], 46, WHITE_C, bold=True, first=True)
    tf = textbox(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.7))
    para(tf, m["subtitle"], 22, rgb("#C9D2EA"), bold=True, first=True)
    meta = m["meta"]
    ent = "entity" if meta["entities"] == 1 else "entities"
    tf = textbox(s, Inches(0.8), Inches(4.2), Inches(10.5), Inches(1.4))
    para(tf, f"Generated from {meta['custodians']} parsed custodian statement(s) across "
             f"{meta['entities']} {ent} · {meta['positions']} positions.", 14, rgb("#AEB8D6"),
         first=True)
    para(tf, "Every figure is computed by the deterministic engine from the client's own "
             "holdings — nothing is invented.", 14, rgb("#AEB8D6"))
    tf = textbox(s, Inches(0.8), Inches(5.9), Inches(11), Inches(0.5))
    para(tf, f"As of {m['as_of']}", 13, rgb("#C8A24A"), bold=True, first=True)

    # ---- Slide 2 · CIO commentary (optional) ----------------------------- #
    if m.get("narrative"):
        s = prs.slides.add_slide(blank)
        header(s, "Investment Commentary", "Chief Investment Office — Commentary")
        tf = textbox(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(4.6))
        first = True
        for pblock in m["narrative"].split("\n\n"):
            if pblock.strip():
                para(tf, pblock.strip(), 14, NAVY_C, first=first, space=10)
                first = False
        tf = textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7))
        para(tf, "Generated prose grounded strictly in the deterministic figures in the "
                 "following slides — it quotes those figures but does not compute or alter "
                 "them. For discussion only; not investment advice.", 10, rgb("#9AA4C4"),
             first=True)
        footer(s, 2)

    # ---- Slide 2 · position ---------------------------------------------- #
    s = prs.slides.add_slide(blank)
    header(s, "Position", "Current Consolidated Position")
    # metric cards as a 5-wide strip of textboxes
    n = len(m["metrics"])
    cw = Inches(12.1 / n)
    for i, (lab, val) in enumerate(m["metrics"]):
        left = Inches(0.6) + int(cw) * i
        box = s.shapes.add_textbox(left, Inches(2.0), int(cw) - Inches(0.15), Inches(1.5))
        tf = box.text_frame; tf.word_wrap = True
        para(tf, lab.upper(), 10, GOLD_C, bold=True, first=True, space=6)
        para(tf, val, 24, NAVY_C, bold=True)
    tf = textbox(s, Inches(0.6), Inches(3.9), Inches(12.1), Inches(2.4))
    para(tf, f"Mandate {meta['mandate']} · Risk appetite {meta['risk']} · Ability to take risk "
             f"{meta['ability']}.", 13, NAVY_C, bold=True, first=True, space=8)
    para(tf, "Consolidated across " + ", ".join(m["custodian_list"]) + ".", 12, SOFT_C)
    para(tf, m["prov"], 11, SOFT_C)
    footer(s, 2 + off)

    # ---- Slide 3 · allocation -------------------------------------------- #
    s = prs.slides.add_slide(blank)
    header(s, "Allocation", "Current Allocation vs Target")
    table(s, ["Asset class", "Value (USD)", "% of gross", "Target", "Drift"], m["alloc_rows"],
          Inches(0.6), Inches(1.9), Inches(12.1), col_ratios=[3, 2.4, 1.8, 1.8, 1.8],
          total=["Gross assets", m["gross_str"], "100.0%", "", ""],
          aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])
    tf = textbox(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.8))
    para(tf, "Weights are computed from parsed statement values against the mandate target. "
             "Sleeves with no target (e.g. alternatives, real estate) are flagged on the "
             "suitability slide, not treated as compliant by omission.", 11, SOFT_C, first=True)
    footer(s, 3 + off)

    # ---- Slide 4 · rebalance --------------------------------------------- #
    s = prs.slides.add_slide(blank)
    header(s, "Proposal", "Rebalancing Proposal — Before → After")
    table(s, ["Sleeve", "Before", "Target", "Trade (USD)", "Action", "Note"], m["reb_rows"],
          Inches(0.6), Inches(1.9), Inches(12.1), col_ratios=[2.3, 1.3, 1.3, 1.8, 1.3, 4.2],
          aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT,
                  PP_ALIGN.LEFT, PP_ALIGN.LEFT], fs=10)
    rs = m["reb_summary"]
    fund = "nets to ≈$0 — self-funding" if rs["selffund"] else "requires external cash / raises proceeds"
    tf = textbox(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.85))
    para(tf, f"Buys {rs['buys']} · Sells {rs['sells']} · Net of trades {rs['net']} ({fund}). "
             f"Illiquid sleeves cannot be traded on demand — amend the mandate to add them to "
             f"target, or stage the reduction against redemption windows.", 11, SOFT_C, first=True)
    footer(s, 4 + off)

    # ---- Slide 5 · suitability ------------------------------------------- #
    s = prs.slides.add_slide(blank)
    header(s, "Suitability", "Suitability of the Proposed Book")
    tf = textbox(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(4.9))
    para(tf, f"Mandate gate · {m['gate']}", 13, NAVY_C, bold=True, first=True, space=10)
    if m["suit_items"]:
        for e, d in m["suit_items"]:
            para(tf, f"[{e.upper()}]  {d}", 12, rgb(ENFORCE.get(e, SOFT)), space=7)
    else:
        para(tf, "Book is within all defined suitability bands.", 12, rgb(OK))
    footer(s, 5 + off)

    # ---- Slide 6 · data & method ----------------------------------------- #
    s = prs.slides.add_slide(blank)
    header(s, "Data & Method", "Data Quality, Analyst Notes & Provenance")
    tf = textbox(s, Inches(0.6), Inches(1.9), Inches(6.0), Inches(4.8))
    para(tf, "Reconciliation & data-quality flags", 14, NAVY_C, bold=True, first=True, space=8)
    if m["data_quality"]:
        for k, d in m["data_quality"]:
            para(tf, f"• [{k}] {d}", 11, SOFT_C, space=6)
    else:
        para(tf, "No data-quality issues detected.", 11, SOFT_C)
    tf = textbox(s, Inches(6.9), Inches(1.9), Inches(5.8), Inches(4.8))
    para(tf, "Analyst notes folded into this proposal", 14, NAVY_C, bold=True, first=True, space=8)
    if m["overlays"]:
        para(tf, "Overlay sleeves: " + " · ".join(m["overlays"]), 11, NAVY_C, bold=True, space=6)
    if m["analyst_notes"]:
        for note in m["analyst_notes"]:
            para(tf, f"• {note}", 11, SOFT_C, space=6)
    elif not m["overlays"]:
        para(tf, "No analyst notes captured.", 11, SOFT_C)
    para(tf, m["prov"] + " For discussion only; not investment advice.", 10, rgb("#9AA4C4"),
         space=6)
    footer(s, 6 + off)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# 3) PDF renderer
# --------------------------------------------------------------------------- #
def render_pdf(m: dict) -> bytes:
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib.units import mm
    from reportlab.lib.colors import HexColor
    from reportlab.lib.utils import simpleSplit
    from reportlab.pdfgen import canvas
    from reportlab.platypus import Table, TableStyle

    PAGE = landscape(A4)  # (841.89, 595.27) pts
    PW, PH = PAGE
    MARGIN = 40
    NAVY_C, GOLD_C, SOFT_C, RULE_C = (HexColor(NAVY), HexColor(GOLD), HexColor(SOFT), HexColor(RULE))
    WHITE = HexColor("#FFFFFF")
    off = 1 if m.get("narrative") else 0

    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=PAGE)

    def header(eyebrow, title):
        c.setFillColor(GOLD_C); c.setFont("Helvetica-Bold", 10)
        c.drawString(MARGIN, PH - MARGIN - 4, eyebrow.upper())
        c.setFillColor(NAVY_C); c.setFont("Helvetica-Bold", 23)
        c.drawString(MARGIN, PH - MARGIN - 32, title)

    def footer(n):
        c.setFillColor(HexColor("#9AA4C4")); c.setFont("Helvetica", 8)
        c.drawString(MARGIN, 22, "Portfolio Proposal · Confidential · computed by the "
                                 "deterministic engine")
        c.drawRightString(PW - MARGIN, 22, str(n))

    def wrapped(text, x, top, width, size=10, color=SOFT_C, leading=13, font="Helvetica"):
        c.setFillColor(color); c.setFont(font, size)
        y = top
        for line in simpleSplit(text, font, size, width):
            c.drawString(x, y, line); y -= leading
        return y

    def draw_table(headers, rows, x, top, width, col_ratios, aligns, total=None, fs=9.5):
        data = [headers] + rows + ([total] if total else [])
        tot = sum(col_ratios)
        col_w = [width * r / tot for r in col_ratios]
        # wrap long text cells into Paragraph-like via reportlab Table auto-height
        from reportlab.platypus import Paragraph
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        base = ParagraphStyle("c", fontName="Helvetica", fontSize=fs, leading=fs + 2,
                              textColor=NAVY_C)
        wrapped_rows = []
        for ri, row in enumerate(data):
            cells = []
            for ci, val in enumerate(row):
                st = ParagraphStyle(f"s{ri}{ci}", parent=base,
                                    alignment={"L": 0, "R": 2, "C": 1}[aligns[ci]],
                                    textColor=WHITE if ri == 0 else NAVY_C,
                                    fontName="Helvetica-Bold" if ri == 0 else "Helvetica")
                cells.append(Paragraph(str(val), st))
            wrapped_rows.append(cells)
        t = Table(wrapped_rows, colWidths=col_w)
        style = [("BACKGROUND", (0, 0), (-1, 0), NAVY_C),
                 ("LINEBELOW", (0, 0), (-1, -1), 0.5, RULE_C),
                 ("TOPPADDING", (0, 0), (-1, -1), 5),
                 ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                 ("LEFTPADDING", (0, 0), (-1, -1), 6),
                 ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                 ("VALIGN", (0, 0), (-1, -1), "MIDDLE")]
        if total:
            style += [("BACKGROUND", (0, -1), (-1, -1), HexColor("#F4F6FB")),
                      ("LINEABOVE", (0, -1), (-1, -1), 1.2, GOLD_C)]
        t.setStyle(TableStyle(style))
        tw, th = t.wrapOn(c, width, PH)
        t.drawOn(c, x, top - th)
        return top - th

    # ---- Slide 1 · cover -------------------------------------------------- #
    c.setFillColor(NAVY_C); c.rect(0, 0, PW, PH, fill=1, stroke=0)
    c.setFillColor(HexColor("#C8A24A")); c.setFont("Helvetica-Bold", 11)
    c.drawString(MARGIN + 10, PH - 150, "PORTFOLIO PROPOSAL · CONFIDENTIAL")
    c.setFillColor(WHITE); c.setFont("Helvetica-Bold", 40)
    c.drawString(MARGIN + 10, PH - 210, m["title"])
    c.setFillColor(HexColor("#C9D2EA")); c.setFont("Helvetica-Bold", 19)
    c.drawString(MARGIN + 10, PH - 250, m["subtitle"])
    meta = m["meta"]
    ent = "entity" if meta["entities"] == 1 else "entities"
    wrapped(f"Generated from {meta['custodians']} parsed custodian statement(s) across "
            f"{meta['entities']} {ent} · {meta['positions']} positions. Every figure is computed "
            f"by the deterministic engine from the client's own holdings — nothing is invented.",
            MARGIN + 10, PH - 300, PW - 2 * MARGIN - 120, size=13, color=HexColor("#AEB8D6"),
            leading=18)
    c.setFillColor(HexColor("#C8A24A")); c.setFont("Helvetica-Bold", 12)
    c.drawString(MARGIN + 10, 70, f"As of {m['as_of']}")
    c.showPage()

    # ---- Slide 2 · CIO commentary (optional) ----------------------------- #
    if m.get("narrative"):
        header("Investment Commentary", "Chief Investment Office — Commentary")
        y = PH - 120
        for pblock in m["narrative"].split("\n\n"):
            if pblock.strip():
                y = wrapped(pblock.strip(), MARGIN, y, PW - 2 * MARGIN, size=12,
                            color=NAVY_C, leading=17) - 12
        wrapped("Generated prose grounded strictly in the deterministic figures in the "
                "following pages — it quotes those figures but does not compute or alter "
                "them. For discussion only; not investment advice.", MARGIN, y - 6,
                PW - 2 * MARGIN, size=9, color=HexColor("#9AA4C4"), leading=12)
        footer(2); c.showPage()

    # ---- Slide 2 · position ---------------------------------------------- #
    header("Position", "Current Consolidated Position")
    n = len(m["metrics"]); cw = (PW - 2 * MARGIN) / n
    for i, (lab, val) in enumerate(m["metrics"]):
        cx = MARGIN + cw * i
        c.setFillColor(GOLD_C); c.setFont("Helvetica-Bold", 9)
        c.drawString(cx, PH - 160, lab.upper())
        c.setFillColor(NAVY_C); c.setFont("Helvetica-Bold", 22)
        c.drawString(cx, PH - 190, val)
    y = wrapped(f"Mandate {meta['mandate']} · Risk appetite {meta['risk']} · Ability to take "
                f"risk {meta['ability']}.", MARGIN, PH - 250, PW - 2 * MARGIN, size=12,
                color=NAVY_C, leading=16, font="Helvetica-Bold")
    y = wrapped("Consolidated across " + ", ".join(m["custodian_list"]) + ".", MARGIN, y - 6,
                PW - 2 * MARGIN, size=11, color=SOFT_C)
    wrapped(m["prov"], MARGIN, y - 4, PW - 2 * MARGIN, size=10, color=SOFT_C)
    footer(2 + off); c.showPage()

    # ---- Slide 3 · allocation -------------------------------------------- #
    header("Allocation", "Current Allocation vs Target")
    y = draw_table(["Asset class", "Value (USD)", "% of gross", "Target", "Drift"],
                   m["alloc_rows"], MARGIN, PH - 120, PW - 2 * MARGIN,
                   [3, 2.4, 1.8, 1.8, 1.8], ["L", "R", "R", "R", "R"],
                   total=["Gross assets", m["gross_str"], "100.0%", "", ""])
    wrapped("Weights are computed from parsed statement values against the mandate target. "
            "Sleeves with no target (e.g. alternatives, real estate) are flagged on the "
            "suitability slide, not treated as compliant by omission.", MARGIN, y - 20,
            PW - 2 * MARGIN, size=10, color=SOFT_C)
    footer(3 + off); c.showPage()

    # ---- Slide 4 · rebalance --------------------------------------------- #
    header("Proposal", "Rebalancing Proposal — Before → After")
    y = draw_table(["Sleeve", "Before", "Target", "Trade (USD)", "Action", "Note"],
                   m["reb_rows"], MARGIN, PH - 120, PW - 2 * MARGIN,
                   [2.3, 1.3, 1.3, 1.8, 1.3, 4.2], ["L", "R", "R", "R", "L", "L"], fs=9)
    rs = m["reb_summary"]
    fund = "nets to ~$0 — self-funding" if rs["selffund"] else "requires external cash / raises proceeds"
    wrapped(f"Buys {rs['buys']} · Sells {rs['sells']} · Net of trades {rs['net']} ({fund}). "
            f"Illiquid sleeves cannot be traded on demand — amend the mandate to add them to "
            f"target, or stage the reduction against redemption windows.", MARGIN, y - 20,
            PW - 2 * MARGIN, size=10, color=SOFT_C)
    footer(4 + off); c.showPage()

    # ---- Slide 5 · suitability ------------------------------------------- #
    header("Suitability", "Suitability of the Proposed Book")
    c.setFillColor(NAVY_C); c.setFont("Helvetica-Bold", 13)
    c.drawString(MARGIN, PH - 130, f"Mandate gate · {m['gate']}")
    y = PH - 160
    if m["suit_items"]:
        for e, d in m["suit_items"]:
            y = wrapped(f"[{e.upper()}]  {d}", MARGIN, y, PW - 2 * MARGIN, size=11,
                        color=HexColor(ENFORCE.get(e, SOFT)), leading=15) - 6
    else:
        wrapped("Book is within all defined suitability bands.", MARGIN, y, PW - 2 * MARGIN,
                size=11, color=HexColor(OK))
    footer(5 + off); c.showPage()

    # ---- Slide 6 · data & method ----------------------------------------- #
    header("Data & Method", "Data Quality, Analyst Notes & Provenance")
    colw = (PW - 2 * MARGIN - 30) / 2
    c.setFillColor(NAVY_C); c.setFont("Helvetica-Bold", 13)
    c.drawString(MARGIN, PH - 130, "Reconciliation & data-quality flags")
    y = PH - 155
    if m["data_quality"]:
        for k, d in m["data_quality"]:
            y = wrapped(f"• [{k}] {d}", MARGIN, y, colw, size=10, color=SOFT_C, leading=13) - 5
    else:
        wrapped("No data-quality issues detected.", MARGIN, y, colw, size=10, color=SOFT_C)
    rx = MARGIN + colw + 30
    c.setFillColor(NAVY_C); c.setFont("Helvetica-Bold", 13)
    c.drawString(rx, PH - 130, "Analyst notes folded into this proposal")
    y = PH - 155
    if m["overlays"]:
        y = wrapped("Overlay sleeves: " + " · ".join(m["overlays"]), rx, y, colw, size=10,
                    color=NAVY_C, leading=13, font="Helvetica-Bold") - 5
    if m["analyst_notes"]:
        for note in m["analyst_notes"]:
            y = wrapped(f"• {note}", rx, y, colw, size=10, color=SOFT_C, leading=13) - 5
    elif not m["overlays"]:
        y = wrapped("No analyst notes captured.", rx, y, colw, size=10, color=SOFT_C) - 5
    wrapped(m["prov"] + " For discussion only; not investment advice.", rx, y - 6, colw,
            size=9, color=HexColor("#9AA4C4"), leading=12)
    footer(6 + off); c.showPage()

    c.save()
    return buf.getvalue()
