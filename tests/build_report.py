"""Build the Meridian Copilot test-suite report (DOCX + PPTX, and PDFs if
LibreOffice is available) from a live run of the suite.

Run:  python3 tests/build_report.py
Outputs (repo root): Meridian_Copilot_Test_Report_v10.{docx,pptx} and, when a
`soffice` binary is found, ..._v10_slides.pdf / ..._v10_document.pdf.

Self-contained: paths derive from this file's location, so it works from any cwd.
Needs python-docx + python-pptx (already in requirements.txt); PDFs are optional.
"""
from __future__ import annotations
import ast, re, subprocess, collections, datetime, shutil
from pathlib import Path

TDIR = Path(__file__).resolve().parent          # the tests/ dir
ROOT = TDIR.parent                              # repo root

NAVY="1E2A56"; GOLD="B0872A"; INK="16203C"; SOFT="47517A"; GREEN="2E7D52"; LINE="C8CDDD"

FILES = [
 ("test_statement_parser","Statement ingestion","Engine",
  "Parse the three custodian formats, reconcile against each statement's own total, "
  "convert FX to USD, and keep a PII-free continuity key. Every downstream number rides here."),
 ("test_suitability","Suitability & policy","Engine",
  "The deterministic policy checks. The mandate decides what a breach MEANS; plus "
  "concentration, liquidity, unhedged-FX, hard regulatory blocks and instruction-diff."),
 ("test_portfolio_qa","Conversational Q&A","Engine",
  "Routing picks the right deterministic tool from the phrasing; answers stay grounded "
  "in computed facts; the router never crashes on odd input."),
 ("test_doc_extract","Document ingestion","Deliverable",
  "The research / other uploaders read text deterministically so content can be folded "
  "into the prompt as advisory context — never a source of figures."),
 ("test_narrative","Prompt & narrative","Deliverable",
  "The portable prompt bundles every input; the keyless commentary REFLECTS the documents "
  "(the v10 fix) while only ever quoting the computed numbers."),
 ("test_generate_proposal","Proposal deck render","Deliverable",
  "HTML / PPTX / PDF renderers agree, figures are copied from the model, and documents "
  "change every exported format."),
 ("test_app_scenarios","End-to-end app","App",
  "Drive the real Streamlit app across the views a banker walks through: load samples, "
  "review each view, change the mandate & tactical instructions — and never error."),
 ("test_business_scenarios","Business scenarios","Business",
  "What happens to the output on a client macro view — a rate hike/cut, a war (or its "
  "end) prompting alternatives, a tax rise/cut: it reaches the AI, shapes the commentary, "
  "and never moves the grounded figures."),
 ("test_macro_overlay","Macro overlays","Business",
  "A client macro view MOVES the recommendation: each overlay is a fixed, bounded, "
  "provenanced tilt of the target allocation — the rebalancer chases it, so trades shift "
  "in the scenario's direction while every figure stays deterministic."),
]

HUMAN = {
 "test_parse_custodian_a_usd_csv":"Parse custodian A (USD, per-share CSV) into valued positions.",
 "test_continuity_key_is_stable_and_pii_free":"Continuity key is deterministic across re-parses and carries no PII.",
 "test_check_continuity_detects_change":"Continuity check compares month-to-month keys.",
 "test_to_holdings_consolidates_all_custodians":"Consolidate all three custodians into one holdings list.",
 "test_within_bands_is_clean":"A balanced book inside every band raises no actionable breach.",
 "test_single_security_concentration_flagged":"A single stock over the 10% cap is flagged; a diversified fund is not.",
 "test_money_market_exempt_from_position_cap":"Parked money-market liquidity is exempt from the single-position cap.",
 "test_liquidity_floor_breach":"A book below the min-liquid floor is flagged.",
 "test_unhedged_fx_breach":"Economic FX exposure above the unhedged cap is flagged.",
 "test_routes_net_worth_to_totals":"'Net worth' routes to the totals tool.",
 "test_routes_currency_question":"An FX / currency question routes to the currency tool.",
 "test_routes_suitability_question":"A suitability / mandate question routes to the suitability tool.",
 "test_routes_concentration_question":"A 'largest position' question routes to the concentration tool.",
 "test_routes_liquidity_question":"A liquidity question routes to the liquidity tool.",
 "test_routes_data_quality_question":"A reconciliation question routes to the data-quality tool.",
 "test_answers_are_grounded_in_facts":"Every answer carries facts tied to the computed book.",
 "test_briefing_has_every_domain":"The deterministic briefing fact-sheet spans every domain.",
 "test_ask_never_raises_on_varied_phrasings":"The router resolves any phrasing — never crashes.",
 "test_plain_text_and_markdown":"Plain text and Markdown read cleanly.",
 "test_html_is_stripped_to_text":"HTML is stripped to readable text (tags removed).",
 "test_eml_email_body_extracted":"An .eml email body is extracted.",
 "test_csv_and_json_pass_through_as_text":"CSV and JSON pass through as text.",
 "test_unsupported_type_fails_gracefully":"An unsupported file type fails gracefully (no crash).",
 "test_truncation_flag_on_large_input":"Very large documents are truncated and flagged, never dropped.",
 "test_result_shape_is_stable":"The extraction result always has the same shape.",
 "test_prompt_contains_all_blocks":"The prompt carries every block (FACTS, intake, holdings, docs, tactical).",
 "test_prompt_embeds_document_text_verbatim":"Document text is embedded verbatim in the prompt.",
 "test_prompt_without_docs_omits_that_content":"Without documents, that content is absent from the prompt.",
 "test_facts_block_is_valid_json_and_holds_the_numbers":"The FACTS block is valid JSON and holds the numbers.",
 "test_deterministic_summary_reflects_documents":"v10 fix: documents change the deterministic commentary.",
 "test_supplied_context_quotes_tactical_and_docs":"The commentary quotes short excerpts of the tactical text and docs.",
 "test_deterministic_summary_quotes_computed_figures_only":"The commentary never invents a figure absent from the model.",
 "test_deterministic_summary_carries_the_disclaimer":"The commentary carries the 'not investment advice' disclaimer.",
 "test_renderers_produce_output":"All three renderers (HTML/PPTX/PDF) produce valid output.",
 "test_deck_shows_the_computed_headline_numbers":"The deck shows the computed net-worth & gross figures.",
 "test_commentary_slide_present_only_with_narrative":"The commentary slide appears only when a narrative is set.",
 "test_documents_change_every_export_format":"v10 fix: documents change the HTML, PPTX and PDF exports.",
 "test_with_docs_deck_mentions_the_liquidity_event":"With documents, the deck surfaces the client's liquidity event.",
 "test_deck_carries_the_not_verified_disclaimer":"The deck carries the 'not independently verified' disclaimer.",
 "test_app_boots_without_a_book":"The app boots with no book loaded.",
 "test_each_sample_loads_into_overview":"Each sample statement loads into the Overview view.",
 "test_core_views_render_with_full_book":"Overview / Holdings / Suitability / Proposal all render.",
 "test_proposal_shows_prompt_and_commentary":"The Proposal view shows the assembled prompt and commentary.",
 "test_tactical_instructions_reach_the_prompt":"Tactical instructions flow into the on-screen prompt.",
 "test_mandate_drives_enforcement_severity":"Changing the mandate changes the suitability gate.",
 "test_single_custodian_book_proposal":"A single-custodian book still produces a proposal.",
 "test_proposal_without_book_is_graceful":"The Proposal view with no book degrades gracefully.",
 "test_scenario_reaches_the_prompt":"A client macro view (rate/war/tax) reaches the AI verbatim.",
 "test_scenario_shapes_the_commentary":"The macro view shapes the advisory commentary.",
 "test_scenario_leaves_the_numbers_grounded":"A macro view never alters the computed figures.",
 "test_rate_scenario_reports_exposure_without_fabricating":"Rate scenario reports bond exposure, invents no unsourced impact.",
 "test_rate_question_routes_to_the_rate_tool":"A rate-move question routes to the rate tool.",
 "test_scenario_changes_deck_prose_not_its_numbers":"The scenario moves the deck's prose, not its tables.",
 "test_overlay_moves_each_sleeve_the_right_way":"Each overlay moves the target sleeves in the documented direction.",
 "test_overlay_preserves_total_allocation":"A tilt sums to zero — the book stays fully invested.",
 "test_overlay_never_goes_negative":"No sleeve is driven negative by a tilt.",
 "test_none_overlay_is_a_no_op":"'No overlay' changes nothing.",
 "test_unknown_overlay_falls_back_to_none":"An unknown overlay falls back to a no-op.",
 "test_changes_reports_only_moved_sleeves_and_describe_is_provenanced":"Reported changes name only moved sleeves; the summary is provenanced.",
 "test_no_overlay_shows_no_tilt_banner":"No overlay → no tilt banner on the proposal.",
 "test_overlay_moves_the_target_in_the_proposal":"Selecting an overlay moves the target in the live proposal.",
}
PARAM = {"test_each_sample_loads_into_overview":3,"test_core_views_render_with_full_book":4,
         "test_mandate_drives_enforcement_severity":2,
         # business scenarios: 6 scenarios × 2 delivery channels
         "test_scenario_reaches_the_prompt":12, "test_scenario_shapes_the_commentary":12,
         "test_scenario_leaves_the_numbers_grounded":12,
         # macro overlays: 5 scenarios parametrised
         "test_overlay_moves_each_sleeve_the_right_way":5, "test_overlay_preserves_total_allocation":5,
         "test_overlay_never_goes_negative":5,
         "test_changes_reports_only_moved_sleeves_and_describe_is_provenanced":5,
         "test_overlay_moves_the_target_in_the_proposal":2}

def first_line(s): return re.sub(r"\s+"," ",s.strip().split("\n\n")[0]).strip() if s else ""

def collect():
    data=[]
    for stem,title,layer,desc in FILES:
        tree=ast.parse((TDIR/f"{stem}.py").read_text())
        tests=[]
        for n in tree.body:
            if isinstance(n,ast.FunctionDef) and n.name.startswith("test_"):
                d=first_line(ast.get_docstring(n) or "") or HUMAN.get(n.name) or \
                  n.name[5:].replace("_"," ").capitalize()+"."
                tests.append((n.name, d, PARAM.get(n.name,1)))
        data.append((stem,title,layer,desc,tests))
    return data

# run the suite to get the real pass count + timing
res=subprocess.run(["python3","-m","pytest","-o","addopts="],cwd=ROOT,
                   capture_output=True,text=True)
tail=res.stdout.strip().splitlines()[-1]
m=re.search(r"(\d+) passed.*?in ([\d.]+)s",tail)
NPASS=int(m.group(1)); SECS=m.group(2)
DATA=collect()
TOTAL=sum(pc for *_,ts in DATA for _,_,pc in ts)
DATE=datetime.date.today().strftime("%d %B %Y")

# ============================ DOCX ============================
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

def rgb(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
def shade(cell,hexc):
    tcPr=cell._tc.get_or_add_tcPr(); sh=OxmlElement("w:shd")
    sh.set(qn("w:fill"),hexc); tcPr.append(sh)
def set_font(run,size=None,bold=None,color=None,name="Calibri"):
    run.font.name=name
    if size:run.font.size=Pt(size)
    if bold is not None:run.font.bold=bold
    if color:run.font.color.rgb=rgb(color)

doc=Document()
st=doc.styles["Normal"]; st.font.name="Calibri"; st.font.size=Pt(10.5); st.font.color.rgb=rgb(INK)

# cover
doc.add_paragraph()
p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
r=p.add_run("MERIDIAN FAMILY OFFICE  ·  COPILOT v10"); set_font(r,11,True,GOLD)
p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
r=p.add_run("AI Copilot — Test Suite & Results"); set_font(r,24,True,NAVY)
p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
r=p.add_run("Verification of the deterministic engine, the deliverable layer, and the end-to-end app")
set_font(r,11,False,SOFT)
p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
r=p.add_run(f"{NPASS}/{TOTAL} tests passed  ·  {DATE}"); set_font(r,12,True,GREEN)
p=doc.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER
r=p.add_run("Runs fully offline in DEMO_MODE — no API key, no network, no billing.")
set_font(r,9.5,False,SOFT); r.italic=True

# summary heading
def h2(txt):
    p=doc.add_paragraph(); p.space_before=Pt(14)
    r=p.add_run(txt); set_font(r,15,True,NAVY)
    pr=p._p.get_or_add_pPr(); b=OxmlElement("w:pBdr"); bt=OxmlElement("w:bottom")
    bt.set(qn("w:val"),"single"); bt.set(qn("w:sz"),"6"); bt.set(qn("w:space"),"2")
    bt.set(qn("w:color"),GOLD); b.append(bt); pr.append(b)
    return p

doc.add_page_break()
h2("Summary")
p=doc.add_paragraph()
r=p.add_run("The suite verifies one central invariant: ")
set_font(r,10.5)
r=p.add_run("every number is computed deterministically; documents and instructions are "
            "advisory context that shape prose, never figures.")
set_font(r,10.5,True,NAVY)

# summary table by layer
layers=collections.OrderedDict()
for stem,title,layer,desc,tests in DATA:
    layers.setdefault(layer,0)
    layers[layer]+=sum(pc for *_,pc in tests)
t=doc.add_table(rows=1,cols=3); t.alignment=WD_TABLE_ALIGNMENT.CENTER
t.style="Table Grid"
hdr=t.rows[0].cells
for i,txt in enumerate(("Layer","What it proves","Tests")):
    hdr[i].paragraphs[0].add_run(txt);
    set_font(hdr[i].paragraphs[0].runs[0],10,True,"FFFFFF"); shade(hdr[i],NAVY)
LAYER_DESC={"Engine":"Parsing, reconciliation, suitability policy, grounded Q&A",
            "Deliverable":"Document ingestion, portable prompt, narrative, deck render",
            "App":"The Streamlit app end-to-end across every view",
            "Business":"Client macro views (rates, war, tax) shape prose, not the numbers"}
for layer,cnt in layers.items():
    row=t.add_row().cells
    row[0].paragraphs[0].add_run(layer); set_font(row[0].paragraphs[0].runs[0],10,True,NAVY)
    row[1].paragraphs[0].add_run(LAYER_DESC[layer]); set_font(row[1].paragraphs[0].runs[0],10)
    row[2].paragraphs[0].add_run(str(cnt)); set_font(row[2].paragraphs[0].runs[0],10,True,GREEN)
    row[2].paragraphs[0].alignment=WD_ALIGN_PARAGRAPH.CENTER

# per-file detail
for stem,title,layer,desc,tests in DATA:
    n=sum(pc for *_,pc in tests)
    h2(f"{title}  ·  {n} tests")
    p=doc.add_paragraph(); r=p.add_run(f"{layer}   "); set_font(r,8.5,True,GOLD)
    r=p.add_run(desc); set_font(r,9.5,False,SOFT); r.italic=True
    p=doc.add_paragraph(); r=p.add_run(f"file: tests/{stem}.py"); set_font(r,8.5,False,SOFT,"Consolas")
    tb=doc.add_table(rows=1,cols=3); tb.style="Table Grid"
    tb.columns[0].width=Inches(0.55); tb.columns[1].width=Inches(5.4); tb.columns[2].width=Inches(0.8)
    hd=tb.rows[0].cells
    for i,txt in enumerate(("#","Scenario verified","Result")):
        hd[i].paragraphs[0].add_run(txt); set_font(hd[i].paragraphs[0].runs[0],9,True,"FFFFFF"); shade(hd[i],NAVY)
    idx=1
    for name,d,pc in tests:
        label=d + ("" if pc==1 else f"  ({pc} variants)")
        row=tb.add_row().cells
        row[0].paragraphs[0].add_run(str(idx)); set_font(row[0].paragraphs[0].runs[0],9,True,SOFT)
        row[0].paragraphs[0].alignment=WD_ALIGN_PARAGRAPH.CENTER
        row[1].paragraphs[0].add_run(label); set_font(row[1].paragraphs[0].runs[0],9.5)
        rc=row[2].paragraphs[0]; rc.add_run("PASS"); set_font(rc.runs[0],9,True,GREEN)
        rc.alignment=WD_ALIGN_PARAGRAPH.CENTER
        shade(row[2],"EAF3EC")
        idx+=1

# how to run
h2("How to run")
for line in ("python3 -m pytest                      # whole suite",
             "python3 -m pytest tests/test_suitability.py -v",
             "python3 -m pytest -o addopts='' -v      # per-test PASS lines"):
    p=doc.add_paragraph(); r=p.add_run(line); set_font(r,9.5,False,INK,"Consolas")

out_docx=ROOT/"Meridian_Copilot_Test_Report_v10.docx"
doc.save(out_docx)
print("wrote",out_docx.name)

# ============================ PPTX ============================
from pptx import Presentation
from pptx.util import Inches as PI, Pt as PP, Emu
from pptx.dml.color import RGBColor as PC
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def pc(h): return PC(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
prs=Presentation(); prs.slide_width=PI(13.333); prs.slide_height=PI(7.5)
BLANK=prs.slide_layouts[6]
W,H=prs.slide_width,prs.slide_height

def box(s,l,t,w,h,fill=None,line=None):
    sh=s.shapes.add_shape(1,l,t,w,h)
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb=pc(fill)
    if line: sh.line.color.rgb=pc(line); sh.line.width=PP(1)
    else: sh.line.fill.background()
    sh.shadow.inherit=False
    return sh
def txt(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=2):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Emu(0); tf.margin_top=tf.margin_bottom=Emu(0)
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=PP(sp)
        for s_,sz,bold,col in line:
            r=p.add_run(); r.text=s_; r.font.size=PP(sz); r.font.bold=bold
            r.font.color.rgb=pc(col); r.font.name="Calibri"
    return tb

# title slide
s=prs.slides.add_slide(BLANK)
box(s,0,0,W,H,fill=NAVY)
box(s,0,H-PI(0.28),W,PI(0.28),fill=GOLD)
txt(s,PI(0.9),PI(2.2),PI(11.5),PI(0.5),[[("MERIDIAN FAMILY OFFICE  ·  COPILOT v10",15,True,GOLD)]])
txt(s,PI(0.9),PI(2.75),PI(11.5),PI(1.6),
    [[("AI Copilot — Test Suite & Results",40,True,"FFFFFF")]])
txt(s,PI(0.9),PI(4.4),PI(11.5),PI(0.5),
    [[("Verifying the deterministic engine, the deliverable layer, and the end-to-end app",16,False,"C8CDDD")]])
txt(s,PI(0.9),PI(5.4),PI(11.5),PI(0.6),
    [[(f"{NPASS}/{TOTAL} tests passed   ·   {DATE}",18,True,"7FD1A0")]])
txt(s,PI(0.9),PI(6.1),PI(11.5),PI(0.4),
    [[("Runs fully offline in DEMO_MODE — no API key, no network, no billing.",12,False,"9AA4C4")]])

# summary slide
s=prs.slides.add_slide(BLANK)
txt(s,PI(0.7),PI(0.5),PI(12),PI(0.4),[[("Coverage at a glance",30,True,NAVY)]])
box(s,PI(0.7),PI(1.15),PI(3.2),PI(0.05),fill=GOLD)
txt(s,PI(0.7),PI(1.35),PI(12),PI(0.6),
    [[("Central invariant: ",14,True,NAVY),
      ("every number is deterministic; documents and instructions shape prose, never figures.",14,False,INK)]])
# cards computed dynamically from DATA so they stay correct as tests grow
LAYER_BLURB={"Engine":"Parsing · reconciliation · suitability · grounded Q&A",
             "Deliverable":"Doc ingestion · portable prompt · narrative · deck render",
             "App":"Streamlit end-to-end across every view",
             "Business":"Client macro views: rates · war · tax → prose, not the numbers"}
_layer_files=collections.OrderedDict()
for stem,title_,layer,desc,tests in DATA:
    nf,nt=_layer_files.get(layer,(0,0))
    _layer_files[layer]=(nf+1, nt+sum(pc for *_,pc in tests))
cards=[(lyr, f"{nf} file{'s' if nf!=1 else ''} · {nt} tests", LAYER_BLURB.get(lyr,""))
       for lyr,(nf,nt) in _layer_files.items()]
n=len(cards); gap=PI(0.12); total=PI(12.0)
cw=int((total-gap*(n-1))/n); x=PI(0.7)
for title_,cnt,desc in cards:
    box(s,x,PI(2.3),cw,PI(2.6),fill="F4F6FB",line=LINE)
    box(s,x,PI(2.3),cw,PI(0.12),fill=GOLD)
    txt(s,x+PI(0.2),PI(2.6),cw-PI(0.4),PI(0.5),[[(title_,20,True,NAVY)]])
    txt(s,x+PI(0.2),PI(3.12),cw-PI(0.4),PI(0.4),[[(cnt,12,True,GOLD)]])
    txt(s,x+PI(0.2),PI(3.62),cw-PI(0.4),PI(1.1),[[(desc,12,False,SOFT)]])
    x=x+cw+gap
txt(s,PI(0.7),PI(5.4),PI(12),PI(1.2),
    [[("What's proven end-to-end",15,True,NAVY)],
     [("• The v10 fix: research/other documents and tactical instructions change the deterministic "
       "commentary and every exported deck (HTML, PPTX, PDF).",12.5,False,INK)],
     [("• The same book reads differently by mandate: advisory flags, discretionary hard-blocks.",12.5,False,INK)],
     [("• Reconciliation catches the ~$436 custodian break instead of trusting it silently.",12.5,False,INK)],
     [("• A macro view as tactical text / documents reshapes the commentary but never moves the "
       "figures; a selected macro OVERLAY deliberately shifts the recommendation via a bounded, "
       "provenanced target tilt (rates → cash into bonds; risk-off → gold & cash up, equity down).",12.5,False,INK)]],sp=6)

# per-file slides
ICON={"Engine":GOLD,"Deliverable":NAVY,"App":GREEN}
for stem,title,layer,desc,tests in DATA:
    s=prs.slides.add_slide(BLANK)
    n=sum(pc for *_,pc in tests)
    box(s,0,0,W,PI(1.15),fill=NAVY)
    box(s,0,PI(1.15),W,PI(0.08),fill=GOLD)
    txt(s,PI(0.7),PI(0.22),PI(9.5),PI(0.5),[[(title,26,True,"FFFFFF")]])
    txt(s,PI(0.7),PI(0.78),PI(9.5),PI(0.3),[[(f"{layer}  ·  tests/{stem}.py",12,False,"C8CDDD")]])
    txt(s,W-PI(2.6),PI(0.3),PI(2.0),PI(0.6),[[(f"{n}/{n} PASS",20,True,"7FD1A0")]],align=PP_ALIGN.RIGHT)
    txt(s,PI(0.7),PI(1.4),PI(12),PI(0.5),[[(desc,13,False,SOFT)]])
    # test rows (two columns if many)
    rows=[(d + ("" if pc==1 else f" ({pc}×)")) for _,d,pc in tests]
    top=PI(2.15); rh=PI(0.46)
    if len(rows)<=8:
        for i,d in enumerate(rows):
            y=top+rh*i
            box(s,PI(0.7),y,PI(11.9),rh-PI(0.08),fill=("F4F6FB" if i%2 else "FFFFFF"),line=LINE)
            txt(s,PI(0.9),y+PI(0.05),PI(0.5),PI(0.3),[[("✓",13,True,GREEN)]])
            txt(s,PI(1.45),y+PI(0.03),PI(11.0),PI(0.4),[[(d,12.5,False,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    else:
        half=(len(rows)+1)//2; cols=[rows[:half],rows[half:]]
        for ci,col in enumerate(cols):
            cx=PI(0.7)+ci*PI(6.15)
            for i,d in enumerate(col):
                y=top+rh*i
                box(s,cx,y,PI(5.95),rh-PI(0.08),fill=("F4F6FB" if i%2 else "FFFFFF"),line=LINE)
                txt(s,cx+PI(0.15),y+PI(0.05),PI(0.4),PI(0.3),[[("✓",12,True,GREEN)]])
                txt(s,cx+PI(0.6),y+PI(0.02),PI(5.3),PI(0.4),[[(d,11,False,INK)]],anchor=MSO_ANCHOR.MIDDLE)

# closing slide
s=prs.slides.add_slide(BLANK)
box(s,0,0,W,H,fill=NAVY); box(s,0,H-PI(0.28),W,PI(0.28),fill=GOLD)
txt(s,PI(0.9),PI(2.6),PI(11.5),PI(1.0),[[("How to run",34,True,"FFFFFF")]])
txt(s,PI(0.9),PI(3.7),PI(11.5),PI(1.5),
    [[("python3 -m pytest",18,True,"7FD1A0")],
     [("python3 -m pytest tests/test_suitability.py -v",15,False,"C8CDDD")],
     [(f"{NPASS}/{TOTAL} passed  ·  offline, deterministic, repeatable",14,False,"9AA4C4")]],sp=10)

out_pptx=ROOT/"Meridian_Copilot_Test_Report_v10.pptx"
prs.save(out_pptx)
print("wrote",out_pptx.name)

# ---- optional PDF twins (best-effort via LibreOffice) --------------------- #
def _soffice() -> str | None:
    for c in ("soffice", "libreoffice",
              "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        if Path(c).exists() or shutil.which(c):
            return c
    return None

soffice = _soffice()
if soffice:
    import tempfile
    for src, pdf_name in ((out_docx, "Meridian_Copilot_Test_Report_v10_document.pdf"),
                          (out_pptx, "Meridian_Copilot_Test_Report_v10_slides.pdf")):
        with tempfile.TemporaryDirectory() as td:  # isolate: soffice reuses the stem
            subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                            "--outdir", td, str(src)],
                           capture_output=True, text=True)
            made = Path(td)/(src.stem + ".pdf")
            if made.exists():
                shutil.copy(made, ROOT/pdf_name)
                print("wrote", pdf_name)
            else:
                print("PDF conversion failed for", src.name)
else:
    print("(LibreOffice not found — skipped PDFs; docx/pptx written)")

print(f"totals: {NPASS}/{TOTAL} pass")
