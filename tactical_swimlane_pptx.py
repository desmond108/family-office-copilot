"""tactical_swimlane_pptx.py — the tactical-instructions process as a native
16:9 PowerPoint swimlane, matching tactical_swimlane_note.html (v7).

Two slides:
  1. The answer — does "Confirm items" fill Allocations & Limits? (the two-click
     reality: Confirm fills nothing; a separate Apply fills the allocation sleeves;
     limits are always manual).
  2. The swimlane — three actor lanes (Client / You / Copilot) across four phases,
     numbered action boxes 1-9, a needs-clarification decision, and the copilot's
     Proposed-allocation step.

House navy / gold, matching the other decks. Run:
    python tactical_swimlane_pptx.py   ->  Tactical_Swimlane.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1E, 0x2A, 0x56); NAVY_DEEP = RGBColor(0x13, 0x1D, 0x3D)
GOLD = RGBColor(0xB0, 0x87, 0x2A); GOLD_DK = RGBColor(0x9C, 0x74, 0x22)
GOLD_TINT = RGBColor(0xF7, 0xEF, 0xDC)
GREEN = RGBColor(0x2F, 0x6F, 0x57); GREEN_TINT = RGBColor(0xE6, 0xF1, 0xEB)
AMBER = RGBColor(0xB0, 0x6A, 0x17); AMBER_TINT = RGBColor(0xFB, 0xE6, 0xD2)
INK = RGBColor(0x16, 0x20, 0x3C); SOFT = RGBColor(0x47, 0x51, 0x7A)
FAINT = RGBColor(0x83, 0x8C, 0xAD); LINE = RGBColor(0xDF, 0xE2, 0xEE)
LINE_STRONG = RGBColor(0xC8, 0xCD, 0xDD); SLATE_TINT = RGBColor(0xEE, 0xF0, 0xF6)
CLIENT_TINT = RGBColor(0xEE, 0xF2, 0xFB)
PAPER = RGBColor(0xFF, 0xFF, 0xFF); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xB9, 0xC2, 0xE0)

SERIF = "Georgia"; SANS = "Calibri"; MONO = "Consolas"

prs = Presentation()
prs.slide_width = In(13.333); prs.slide_height = In(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08, dash=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        if dash:  # solidFill is already set above, so prstDash is in schema order
            from pptx.oxml.ns import qn
            ln = shp.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = In(0.08); tf.margin_right = In(0.07)
    tf.margin_top = In(0.05); tf.margin_bottom = In(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    try:
        tf.auto_size = None
    except Exception:
        pass
    return shp


def txt(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    try:
        tf.auto_size = None
    except Exception:
        pass
    return tf


def para(tf, runs, *, first=False, align=PP_ALIGN.LEFT, before=0, after=0, spacing=1.0):
    """runs: list of (text, size, color, bold, font). Single paragraph, multi-run."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(before); p.space_after = Pt(after)
    try:
        p.line_spacing = spacing
    except Exception:
        pass
    for text, size, color, bold, font in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font; r.font.color.rgb = color
    return p


def one(tf, text, size, color, **kw):
    return para(tf, [(text, size, color, kw.pop("bold", False), kw.pop("font", SANS))], **kw)


# =========================================================================== #
# Slide 1 — the two-click answer
# =========================================================================== #
s = slide()
one(txt(s, 0.8, 0.5, 11.5, 0.4), "MERIDIAN FAMILY OFFICE COPILOT · v7", 11.5, GOLD,
    bold=True, first=True)
one(txt(s, 0.8, 0.92, 11.7, 0.9),
    "Does “Confirm items” fill your Allocations & Limits?", 27, NAVY, font=SERIF,
    bold=True, first=True)

cards = [
    ("Confirm items", "commits the interpreted instructions", NAVY,
     "Produces a 📡 monitoring watchlist (triggers) and guidance folded into the "
     "proposal. It writes into NO field.", SLATE_TINT, LINE_STRONG, INK),
    ("Apply to allocation targets", "the second, separate click", GOLD_DK,
     "The ONLY action that fills the allocation sleeves. Appears only if the client "
     "stated target weights (“Gold ETF: 20%”). You still verify and adjust.",
     GOLD_TINT, GOLD, GOLD_DK),
    ("Limits", "band tolerance · liquidity · FX · position caps", SOFT,
     "Always MANUAL — never extracted from the instructions or auto-filled.",
     PAPER, LINE_STRONG, INK),
]
cy = 2.25; ch = 1.25; gap = 0.24
for i, (title, kicker, tcol, body, fill, lcol, bodycol) in enumerate(cards):
    y = cy + i * (ch + gap)
    b = rect(s, 0.8, y, 11.73, ch, fill=fill, line=lcol, line_w=1.5, radius=0.06)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(f"{i+1}.  ", 15, tcol, True, SANS), (title, 15, tcol, True, SANS),
              (f"   — {kicker}", 12, SOFT, False, SANS)], first=True)
    one(tf, body, 12.5, bodycol, before=5, spacing=1.08)

b = rect(s, 0.8, 6.55, 11.73, 0.88, fill=NAVY_DEEP, radius=0.06)
tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, [("Order & Analyse:  ", 12, GOLD, True, SANS),
          ("Loading the statements and setting the mandate / risk / allocations are independent — "
           "do them in any order. Analyse processes the statements; after Confirm the views update "
           "live (no re-Analyse).", 12, WHITE, False, SANS)], first=True, spacing=1.06)


# =========================================================================== #
# Slide 2 — the swimlane
# =========================================================================== #
s = slide()
one(txt(s, 0.5, 0.34, 11.5, 0.4), "MERIDIAN FAMILY OFFICE COPILOT · v7 · PROCESS",
    11, GOLD, bold=True, first=True)
one(txt(s, 0.5, 0.68, 12.4, 0.6), "Tactical instructions — who does what, and when",
    23, NAVY, font=SERIF, bold=True, first=True)

# geometry
LX, LW = 0.4, 0.95
GX = LX + LW                      # 1.35
GR = 13.0
PH_Y, PH_H = 1.5, 0.38
# phase widths (C widest — most action): A, B, C, D
PW = [2.02, 2.5, 4.66, 2.47]
PXS = [GX]
for w in PW:
    PXS.append(PXS[-1] + w)      # PXS[i]..PXS[i+1] is phase i

LANES = [("Client", "Source", CLIENT_TINT, 1.92, 1.18),
         ("You", "Analyst", RGBColor(0xFB, 0xF6, 0xE9), 3.16, 2.46),
         ("Copilot", "System", SLATE_TINT, 5.68, 1.5)]

# phase headers
phase_titles = [("A", "Ingest"), ("B", "Mandate & policy"),
                ("C", "Capture tactical instructions"), ("D", "Review & generate")]
rect(s, LX, PH_Y, LW, PH_H, fill=NAVY_DEEP, radius=0.04)
for i, (pn, pt) in enumerate(phase_titles):
    b = rect(s, PXS[i], PH_Y, PW[i], PH_H, fill=NAVY_DEEP, radius=0.04)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(f"{pn}  ", 11, GOLD, True, MONO), (pt, 11.5, WHITE, True, SANS)],
         first=True, align=PP_ALIGN.CENTER)

# lane bands + labels
for name, sub, bg, y, h in LANES:
    rect(s, GX, y, GR - GX, h, fill=bg, line=LINE, line_w=0.75, radius=0.01)
    lb = rect(s, LX, y, LW, h, fill=None, line=None)
    tf = lb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    col = {"You": GOLD_DK, "Client": NAVY, "Copilot": SOFT}[name]
    para(tf, [(name, 13, col, True, SANS)], first=True, align=PP_ALIGN.CENTER)
    para(tf, [(sub.upper(), 8.5, FAINT, False, SANS)], align=PP_ALIGN.CENTER, before=1)

CY = {name: (y, h) for name, _, _, y, h in LANES}


def abox(px, y, w, h, kind, num, title, desc, tsize=10, dsize=8.3):
    fill, lcol, lw, dash, ncol = {
        "you": (PAPER, GOLD, 1.25, False, GOLD_DK),
        "bot": (PAPER, LINE_STRONG, 1.0, True, FAINT),
        "cli": (PAPER, RGBColor(0x9C, 0xA8, 0xCE), 1.0, False, NAVY),
        "dec": (AMBER_TINT, AMBER, 1.25, False, AMBER),
    }[kind]
    b = rect(s, px, y, w, h, fill=fill, line=lcol, line_w=lw, radius=0.10, dash=dash)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    head = []
    if num:
        head.append((f"{num} · ", tsize, ncol, True, SANS))
    head.append((title, tsize, INK if kind != "bot" else SOFT, True, SANS))
    para(tf, head, first=True, spacing=1.0)
    if desc:
        one(tf, desc, dsize, SOFT, before=2, spacing=1.02)
    return b


def gear(px, y, w, h, tag, title, desc):
    b = rect(s, px, y, w, h, fill=PAPER, line=LINE_STRONG, line_w=1.0, radius=0.10, dash=True)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    one(tf, tag, 8, FAINT, bold=True, first=True)
    one(tf, title, 10, SOFT, bold=True, before=2, spacing=1.0)
    if desc:
        one(tf, desc, 8.3, SOFT, before=2, spacing=1.0)
    return b

# ---- CLIENT lane ---- #
cy, ch = CY["Client"]
abox(PXS[0] + 0.12, cy + 0.2, PW[0] - 0.24, ch - 0.4, "cli", None,
     "Hands over statements", "INPUT", tsize=10.5, dsize=8)
abox(PXS[2] + 0.12, cy + 0.2, PW[2] * 0.62, ch - 0.4, "cli", None,
     "Gives instructions (plain language)", "“gold below $4,000”, “buy in tranches”…",
     tsize=10.5, dsize=8)

# ---- YOU lane ---- #
yy, yh = CY["You"]
# Phase A: box 1
abox(PXS[0] + 0.12, yy + 0.75, PW[0] - 0.24, 0.95, "you", "1",
     "Upload → Analyse ▸", "any time · independent of 2–3")
# Phase B: boxes 2 (top) + 3 (bottom)
bx, bw = PXS[1] + 0.12, PW[1] - 0.24
abox(bx, yy + 0.24, bw, 1.02, "you", "2", "Set mandate, risk & ability", None)
abox(bx, yy + 1.34, bw, 1.02, "you", "3", "Set allocation targets & limits",
     "limits manual · sleeves can be pre-filled by 7")
# Phase C: 2x2 grid + decision strip
cx0 = PXS[2] + 0.1; cW = PW[2] - 0.2
colw = (cW - 0.16) / 2
c1, c2 = cx0, cx0 + colw + 0.16
abox(c1, yy + 0.12, colw, 0.72, "you", "4", "Paste → Sort into items", None, tsize=9.5)
abox(c2, yy + 0.12, colw, 0.72, "you", "5", "Review: Keep · edit · notes", None, tsize=9.5)
# decision strip full width
db = rect(s, cx0, yy + 0.9, cW, 0.5, fill=AMBER_TINT, line=AMBER, line_w=1.25, radius=0.12)
tf = db.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, [("◇ needs-clarification?  ", 9.5, AMBER, True, SANS),
          ("Yes → ask client ⟲   ·   No ↓", 9, SOFT, False, SANS)], first=True,
     align=PP_ALIGN.CENTER)
abox(c1, yy + 1.48, colw, 0.86, "you", "6", "Confirm items",
     "commits guidance + watchlist · fills NO fields", tsize=9.5, dsize=7.8)
b7 = abox(c2, yy + 1.48, colw, 0.86, "you", "7", "Apply to allocation targets",
          "only action that fills sleeves · optional · ⤶ fills 3", tsize=9.5, dsize=7.8)
# Phase D: boxes 8 + 9
dx, dw = PXS[3] + 0.12, PW[3] - 0.24
abox(dx, yy + 0.24, dw, 1.02, "you", "8", "Review Overview / Suitability", None)
abox(dx, yy + 1.34, dw, 1.02, "you", "9", "Open Proposal → Generate deck", "PPTX / PDF")

# ---- COPILOT lane ---- #
py, ph = CY["Copilot"]
gear(PXS[0] + 0.12, py + 0.2, PW[0] - 0.24, ph - 0.4, "⚙ AUTO",
     "Parse → build the book", None)
gear(PXS[1] + 0.12, py + 0.2, PW[1] - 0.24, ph - 0.4, "⚙ AUTO · LIVE",
     "Compute allocation, drift, suitability", None)
cpw = (PW[2] - 0.2 - 0.14) / 2
gear(PXS[2] + 0.1, py + 0.2, cpw, ph - 0.4, "⚙ ON CONFIRM",
     "Watchlist + guidance", "unclear items held out")
gear(PXS[2] + 0.1 + cpw + 0.14, py + 0.2, cpw, ph - 0.4, "⚙ IF WEIGHTS GIVEN",
     "Builds Proposed allocation", "you Apply it in 7")
gear(PXS[3] + 0.12, py + 0.2, PW[3] - 0.24, ph - 0.4, "⚙ AUTO",
     "Generate PPTX / PDF", "figures deterministic")

# footer key point
one(txt(s, 0.5, 7.16, 12.5, 0.3),
    "Key point: setup (1–7) is any-order; Confirm (6) fills no fields; allocations only via "
    "Apply (7); limits always manual.", 10, SOFT, bold=True, first=True)

OUT = "Tactical_Swimlane.pptx"
prs.save(OUT)
print(f"wrote {OUT} — {len(prs.slides._sldIdLst)} slides")
