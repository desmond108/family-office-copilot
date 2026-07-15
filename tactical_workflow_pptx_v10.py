"""tactical_workflow_pptx.py — the v10 tactical-instructions workflow as a native
16:9 PowerPoint deck, for dropping into partner / client decks.

Six slides that mirror the shareable artifact (tactical_workflow_note.html): how the
copilot turns the client's tactical instructions + documents into a proposal in v10
— one self-contained prompt handed to the AI Model, which writes the narrative while
the engine computes every number. House navy / gold style, matching generate_proposal.py.

Run:  python tactical_workflow_pptx.py   ->  Tactical_Workflow_v10.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (matches the artifact + generate_proposal) ------------------- #
NAVY = RGBColor(0x1E, 0x2A, 0x56)
NAVY_DEEP = RGBColor(0x13, 0x1D, 0x3D)
GOLD = RGBColor(0xB0, 0x87, 0x2A)
GOLD_DK = RGBColor(0x9C, 0x74, 0x22)
GOLD_TINT = RGBColor(0xF7, 0xEF, 0xDC)
GREEN = RGBColor(0x2F, 0x6F, 0x57)
GREEN_TINT = RGBColor(0xE6, 0xF1, 0xEB)
INK = RGBColor(0x16, 0x20, 0x3C)
SOFT = RGBColor(0x47, 0x51, 0x7A)
FAINT = RGBColor(0x83, 0x8C, 0xAD)
LINE = RGBColor(0xDF, 0xE2, 0xEE)
LINE_STRONG = RGBColor(0xC8, 0xCD, 0xDD)
SLATE_TINT = RGBColor(0xEE, 0xF0, 0xF6)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xB9, 0xC2, 0xE0)

SERIF = "Georgia"          # closest ubiquitous serif to the artifact's Palatino stack
SANS = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = In(13.333)
prs.slide_height = In(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _no_autosize(tf):
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    # corner radius
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = In(0.16); tf.margin_right = In(0.16)
    tf.margin_top = In(0.12); tf.margin_bottom = In(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _no_autosize(tf)
    return shp


def txt(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame; _no_autosize(tf)
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, *, font=SANS, bold=False, italic=False,
         first=False, before=0, after=0, align=PP_ALIGN.LEFT, spacing=1.0, ls=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before); p.space_after = Pt(after)
    try:
        p.line_spacing = spacing
    except Exception:
        pass
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    if ls is not None:
        _set_letter_spacing(r, ls)
    return p


def _set_letter_spacing(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def snum(s, n, light=False):
    tf = txt(s, 11.3, 0.34, 1.7, 0.4)
    para(tf, f"{n:02d} / 06", 10.5, (CREAM if light else FAINT), font=MONO,
         first=True, align=PP_ALIGN.RIGHT, ls=0.4)


def eyebrow(s, text, x=0.86, y=0.72):
    tf = txt(s, x, y, 11, 0.4)
    para(tf, text.upper(), 11.5, GOLD, font=SANS, bold=True, first=True, ls=1.6)


def title(s, text, x=0.86, y=1.12, w=11.6, size=26, color=NAVY):
    tf = txt(s, x, y, w, 1.4)
    para(tf, text, size, color, font=SERIF, bold=True, first=True, spacing=1.05)
    return tf


def sub(s, text, x=0.86, y=2.02, w=11.2, size=14):
    tf = txt(s, x, y, w, 0.7)
    para(tf, text, size, SOFT, first=True, spacing=1.1)


# =========================================================================== #
# Slide 1 — title
# =========================================================================== #
s = slide(NAVY_DEEP)
box(s, 0, 0, 0.14, 7.5, fill=GOLD)                    # gold rail
snum(s, 1, light=True)
eyebrow(s, "The tactical-instructions workflow", x=0.9, y=1.4)
tf = txt(s, 0.9, 2.4, 10.6, 2.4)
para(tf, "From the client's words to a working proposal", 38, WHITE,
     font=SERIF, bold=True, first=True, spacing=1.02)
tf = txt(s, 0.9, 4.7, 9.8, 1.0)
para(tf, "How the copilot turns the client's instructions and documents into a proposal — "
         "and the one rule that keeps every number trustworthy.", 15.5, CREAM, first=True, spacing=1.2)
tf = txt(s, 0.9, 6.5, 11.5, 0.5)
para(tf, "Meridian Family Office Copilot   ·   v10 · Tactical instructions   ·   "
         "For partner & client discussion", 11.5, RGBColor(0x7F, 0x8C, 0xBB), first=True)


# =========================================================================== #
# Slide 2 — pipeline at a glance
# =========================================================================== #
s = slide()
snum(s, 2)
eyebrow(s, "At a glance")
title(s, "Three steps in, one prompt, one proposal")
sub(s, "The first three steps are what the analyst does. Everything to their right is automatic — "
       "and fully visible in the prompt.")

nodes = [
    ("01", "Paste", "Tactical text + documents", SLATE_TINT, INK, LINE),
    ("02", "Set policy", "Mandate · risk · limits", SLATE_TINT, INK, LINE),
    ("03", "Analyse", "Engine computes FACTS", SLATE_TINT, INK, LINE),
    ("→", "Prompt", "One self-contained brief", GOLD_TINT, GOLD_DK, GOLD),
    ("→", "AI Model", "Writes the narrative", GOLD_TINT, GOLD_DK, GOLD),
    ("→", "Deck", "PPTX / PDF", GOLD_TINT, GOLD_DK, GOLD),
]
n = len(nodes); gap = 0.24; x0 = 0.86; total = 11.6
w = (total - gap * (n - 1)) / n
y = 3.1; h = 1.9
for i, (k, t, d, fill, tcol, lcol) in enumerate(nodes):
    x = x0 + i * (w + gap)
    b = box(s, x, y, w, h, fill=fill, line=lcol, line_w=1.25, radius=0.10)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP
    para(tf, k, 12, FAINT if fill == SLATE_TINT else tcol, font=MONO, bold=True, first=True, ls=0.4)
    para(tf, t, 15, tcol, bold=True, before=4, spacing=1.0)
    para(tf, d, 11, SOFT, before=3, spacing=1.05)
    if i < n - 1:
        cx = x + w + (gap - 0.16) / 2
        ar = txt(s, cx - 0.05, y + h / 2 - 0.2, 0.3, 0.4)
        para(ar, "›", 20, LINE_STRONG, first=True, align=PP_ALIGN.CENTER)

tf = txt(s, 0.86, 5.35, 11, 0.5)
para(tf, "The rest of this deck follows what the copilot assembles and what the AI Model does with it.",
     13, SOFT, first=True, bold=True)


# =========================================================================== #
# Slide 3 — what goes into the one prompt
# =========================================================================== #
s = slide()
snum(s, 3)
eyebrow(s, "Step 1 · assemble")
title(s, "One self-contained prompt — everything in one place")
sub(s, "The copilot assembles a single prompt: the deterministic FACTS plus the raw context. "
       "Numbers come only from FACTS; the rest is intent and context.")

blocks = [
    ("ROLE + GROUNDING RULES", "analyst persona · FACTS-only", INK, False),
    ("INTAKE PARAMETERS", "mandate · risk · limits · targets", INK, False),
    ("FACTS (JSON)", "the only source of numbers", GOLD_DK, True),
    ("HOLDINGS + STATEMENT SOURCE", "parsed positions + raw text", INK, False),
    ("RESEARCH / OTHER DOCUMENTS", "full text, as context", INK, False),
    ("TACTICAL INSTRUCTIONS", "the client's words, verbatim", INK, False),
]
bw = (11.6 - 0.28) / 3; bh = 1.18; bx0 = 0.86; bgx = 0.14; bgy = 0.2; by0 = 3.0
for i, (k, v, vcol, hot) in enumerate(blocks):
    col = i % 3; row = i // 3
    x = bx0 + col * (bw + bgx)
    yy = by0 + row * (bh + bgy)
    fill = GOLD_TINT if hot else SLATE_TINT
    lcol = GOLD if hot else LINE
    b = box(s, x, yy, bw, bh, fill=fill, line=lcol, line_w=1.5 if hot else 1.0, radius=0.08)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, k, 11.5, vcol, bold=True, first=True, ls=0.4)
    para(tf, v, 11.5, SOFT, before=6, spacing=1.05)

tf = txt(s, 0.86, 5.9, 11.6, 0.8)
para(tf, "Shown on the Proposal page — editable, copyable, downloadable — so the client can read "
         "exactly what the system asks the AI Model to do.", 12.5, GREEN, bold=True, first=True, spacing=1.1)


# =========================================================================== #
# Slide 4 — the transparency surface
# =========================================================================== #
s = slide()
snum(s, 4)
eyebrow(s, "Step 2 · generate")
title(s, "Read it, edit it, test it anywhere")

sw = 5.68; sh = 3.15; sy = 2.75
# Stream A — generate in-app
b = box(s, 0.86, sy, sw, sh, fill=GOLD_TINT, line=GOLD, line_w=1.5, radius=0.05)
tf = b.text_frame
para(tf, "✨   Generate with the AI Model", 15, INK, bold=True, first=True)
para(tf, "Runs the live AI model on the prompt. Its narrative folds into the proposal deck as a "
         "commentary slide. Falls back to a deterministic summary with no key.", 12.8, SOFT,
     before=8, spacing=1.15)
inner = box(s, 1.12, sy + 1.9, sw - 0.52, 0.95, fill=PAPER, line=LINE_STRONG, line_w=1.0, radius=0.08)
tfi = inner.text_frame; tfi.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfi, "Model-agnostic UI — the button and badge read “AI Model”, never a vendor name.",
     11.5, SOFT, first=True, spacing=1.1)

# Stream B — copy to any AI model
b = box(s, 0.86 + sw + 0.24, sy, sw, sh, fill=SLATE_TINT, line=LINE_STRONG, line_w=1.25, radius=0.05)
tf = b.text_frame
para(tf, "\U0001F4CB   Copy into any AI model", 15, INK, bold=True, first=True)
para(tf, "The prompt is self-contained, so pasting it verbatim into any external AI model "
         "reproduces the proposal — the way the client tests the system across models.", 12.8, SOFT,
     before=8, spacing=1.15)
inner = box(s, 0.86 + sw + 0.24 + 0.26, sy + 1.9, sw - 0.52, 0.95, fill=PAPER,
            line=LINE_STRONG, line_w=1.0, radius=0.08)
tfi = inner.text_frame; tfi.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfi, "Edit before generating; the FACTS block stays the only source of numbers.",
     11.5, SOFT, first=True, spacing=1.1)

# note bar
nb = box(s, 0.86, 6.1, 11.6, 0.75, fill=SLATE_TINT, line=LINE, line_w=1.0, radius=0.06)
tf = nb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "🧠  Both paths use the same prompt. There is no classification, review table, or "
         "enforcement step in between — the analyst reviews the prompt itself.", 12, SOFT,
     first=True, spacing=1.12)


# =========================================================================== #
# Slide 5 — the proposal
# =========================================================================== #
s = slide()
snum(s, 5)
eyebrow(s, "Step 3 · the proposal")
title(s, "AI prose around deterministic tables")

dw = 3.72; dh = 3.55; dy = 2.65; dx0 = 0.86; dgap = 0.22
dests = [
    ("\U0001F9E0", "The prompt", "On the Proposal page — exactly what the AI Model was asked to do.",
     "TRANSPARENCY", GOLD_DK, GOLD, True),
    ("\U0001F4AC", "CIO commentary", "The AI Model's narrative, shaped by the parameters, documents "
     "and tactical instructions — quoting only the computed figures.", "GENERATED PROSE", INK, LINE, False),
    ("\U0001F4CA", "Deterministic tables", "Allocation, drift, rebalance and suitability — computed "
     "by the engine and unchanged by the AI.", "PPTX · PDF", INK, LINE, False),
]
for i, (ic, t, d, where, tcol, lcol, hot) in enumerate(dests):
    x = dx0 + i * (dw + dgap)
    b = box(s, x, dy, dw, dh, fill=PAPER, line=lcol, line_w=1.5 if hot else 1.0, radius=0.06)
    tf = b.text_frame
    para(tf, ic, 22, INK, first=True)
    para(tf, t, 14, tcol, bold=True, before=8)
    para(tf, d, 12, SOFT, before=6, spacing=1.15)
    lab = txt(s, x + 0.16, dy + dh - 0.5, dw - 0.32, 0.35)
    para(lab, where, 10, FAINT, bold=True, first=True, ls=0.6)


# =========================================================================== #
# Slide 6 — guardrail
# =========================================================================== #
s = slide()
snum(s, 6)
eyebrow(s, "The guardrail")
title(s, "The AI writes prose — never a number")

gw = 5.68; gh = 2.45; gy = 2.5
b = box(s, 0.86, gy, gw, gh, fill=GOLD_TINT, line=GOLD, line_w=1.5, radius=0.06)
tf = b.text_frame
para(tf, "THE AI MODEL CAN SHAPE", 11, GOLD_DK, bold=True, first=True, ls=0.8)
for i, li in enumerate(["The CIO commentary and how it reads",
                        "Which considerations to weigh, from the documents & instructions",
                        "The emphasis and ordering of the narrative"]):
    para(tf, "•  " + li, 12.5, SOFT, before=(10 if i == 0 else 6), spacing=1.1)

b = box(s, 0.86 + gw + 0.24, gy, gw, gh, fill=GREEN_TINT, line=GREEN, line_w=1.5, radius=0.06)
tf = b.text_frame
para(tf, "THE AI MODEL NEVER INVENTS", 11, GREEN, bold=True, first=True, ls=0.8)
for i, li in enumerate(["A figure — every number comes from the FACTS block",
                        "Suitability thresholds or the rebalance",
                        "Holdings — every dollar from the real book"]):
    para(tf, "•  " + li, 12.5, SOFT, before=(10 if i == 0 else 6), spacing=1.1)

tf = txt(s, 0.86, gy + gh + 0.45, 11.6, 1.4)
para(tf, "Qualitative claims drawn from the client's documents are context — not independently "
         "verified. Every figure is computed deterministically by the engine.", 18, NAVY, font=SERIF,
     italic=True, bold=True, first=True, spacing=1.15)


OUT = "Tactical_Workflow_v10.pptx"
prs.save(OUT)
print(f"wrote {OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
