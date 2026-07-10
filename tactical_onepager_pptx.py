"""tactical_onepager_pptx.py — the whole "after Confirm" tactical-instructions
workflow condensed onto a SINGLE 16:9 slide, for partners who prefer one diagram
to the six-slide sequence. House navy / gold, matching tactical_workflow_pptx.py.

Reads left → right: the client's words → Sort + Confirm → two outputs (the
monitoring watchlist, and analyst guidance that flows into the deck + commentary)
→ a full-width guardrail band (guidance never moves a computed number).

Run:  python tactical_onepager_pptx.py   ->  Tactical_Workflow_OnePager.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (identical to the six-slide deck) ---------------------------- #
NAVY = RGBColor(0x1E, 0x2A, 0x56)
NAVY_DEEP = RGBColor(0x13, 0x1D, 0x3D)
GOLD = RGBColor(0xB0, 0x87, 0x2A)
GOLD_DK = RGBColor(0x9C, 0x74, 0x22)
GOLD_TINT = RGBColor(0xF7, 0xEF, 0xDC)
GREEN = RGBColor(0x2F, 0x6F, 0x57)
GREEN_TINT = RGBColor(0xE6, 0xF1, 0xEB)
INK = RGBColor(0x16, 0x20, 0x3C)
SOFT = RGBColor(0x47, 0x51, 0x7A)
FAINT = RGBColor(0x83, 0x8C, 0xAD)
LINE = RGBColor(0xDF, 0xE2, 0xEE)
LINE_STRONG = RGBColor(0xC8, 0xCD, 0xDD)
SLATE_TINT = RGBColor(0xEE, 0xF0, 0xF6)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xB9, 0xC2, 0xE0)

SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = In(13.333)
prs.slide_height = In(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = PAPER


def box(x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = In(0.15); tf.margin_right = In(0.15)
    tf.margin_top = In(0.11); tf.margin_bottom = In(0.11)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    try:
        tf.auto_size = None
    except Exception:
        pass
    return shp


def txt(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    try:
        tf.auto_size = None
    except Exception:
        pass
    return tf


def para(tf, text, size, color, *, font=SANS, bold=False, italic=False,
         first=False, before=0, after=0, align=PP_ALIGN.LEFT, spacing=1.0, ls=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before); p.space_after = Pt(after)
    try:
        p.line_spacing = spacing
    except Exception:
        pass
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    if ls is not None:
        r._r.get_or_add_rPr().set("spc", str(int(ls * 100)))
    return p


def chevron(cx, cy):
    t = txt(cx - 0.18, cy - 0.26, 0.5, 0.52, anchor=MSO_ANCHOR.MIDDLE)
    para(t, "›", 30, LINE_STRONG, first=True, align=PP_ALIGN.CENTER)


# ---- header --------------------------------------------------------------- #
tf = txt(0.8, 0.5, 11.5, 0.4)
para(tf, "MERIDIAN FAMILY OFFICE COPILOT · v8", 11.5, GOLD, bold=True, first=True, ls=1.6)
tf = txt(0.8, 0.92, 11.7, 0.9)
para(tf, "The tactical-instructions workflow, at a glance", 27, NAVY, font=SERIF,
     bold=True, first=True)
tf = txt(0.8, 1.72, 11.7, 0.5)
para(tf, "What the copilot does the moment you confirm the sorted items — read left to right.",
     14, SOFT, first=True)

SPINE = 3.75  # vertical centre of the pipeline row

# ---- Stage 1 · client's words --------------------------------------------- #
x1, w1, h1 = 0.8, 3.05, 2.9
b = box(x1, SPINE - h1 / 2, w1, h1, fill=SLATE_TINT, line=LINE_STRONG, line_w=1.25, radius=0.06)
tf = b.text_frame
para(tf, "1 · CLIENT'S WORDS", 11, GOLD_DK, bold=True, first=True, ls=0.8)
para(tf, "In plain language", 11.5, SOFT, before=3)
for i, q in enumerate(['"gold below USD 4,000/oz"',
                       '"add Nasdaq after a 15–20% pullback"',
                       '"buy the bond fund in tranches"',
                       '"low fees, good liquidity"',
                       '"impact of rate hikes?"']):
    para(tf, q, 10.5, INK, font=MONO, before=(11 if i == 0 else 5), spacing=1.05)

# ---- Stage 2 · sort + confirm --------------------------------------------- #
x2, w2, h2 = x1 + w1 + 0.5, 2.25, 2.0
b = box(x2, SPINE - h2 / 2, w2, h2, fill=GREEN_TINT, line=GREEN, line_w=1.5, radius=0.07)
tf = b.text_frame
para(tf, "2 · SORT + CONFIRM", 11, GREEN, bold=True, first=True, ls=0.8)
para(tf, "The copilot types each ask; the analyst reviews and confirms in a table.",
     12, SOFT, before=9, spacing=1.15)
para(tf, "✓ typed & confirmed", 12, GREEN, bold=True, before=11)

chevron(x1 + w1 + 0.25, SPINE)
chevron(x2 + w2 + 0.25, SPINE)

# ---- Stage 3 · three outputs (stacked) ------------------------------------ #
x3 = x2 + w2 + 0.5
w3 = 12.53 - x3
oh = 1.02
gap = 0.22
top = SPINE - (oh * 3 + gap * 2) / 2

# Output A — monitoring watchlist
b = box(x3, top, w3, oh, fill=GOLD_TINT, line=GOLD, line_w=1.5, radius=0.08)
tf = b.text_frame
para(tf, "\U0001F4E1  Monitoring watchlist", 13, INK, bold=True, first=True)
para(tf, "Level-based triggers to watch the book against.", 10.5, SOFT, before=3, spacing=1.04)
p = tf.add_paragraph(); p.space_before = Pt(4)
for seg, col, bold in [("S&P 500 · ", INK, True), ("−15/−20%", GOLD_DK, True),
                       ("     Gold · ", INK, True), ("< $4,000", GOLD_DK, True)]:
    r = p.add_run(); r.text = seg; r.font.size = Pt(10.5); r.font.name = MONO
    r.font.color.rgb = col; r.font.bold = bold
lab = txt(x3 + w3 - 2.4, top + 0.1, 2.28, 0.3, anchor=MSO_ANCHOR.TOP)
para(lab, "LIVES IN THE COPILOT", 9, GOLD_DK, bold=True, first=True, align=PP_ALIGN.RIGHT, ls=0.5)

# Output B — analyst guidance
by = top + oh + gap
b = box(x3, by, w3, oh, fill=SLATE_TINT, line=LINE_STRONG, line_w=1.25, radius=0.08)
tf = b.text_frame
para(tf, "\U0001F9ED  Analyst guidance", 13, INK, bold=True, first=True)
para(tf, "Intent that shapes the written advice — never the figures.", 10.5, SOFT,
     before=3, spacing=1.04)
para(tf, "Execution style · selection criteria · open questions", 10.5, SOFT, font=MONO, before=4)
lab = txt(x3 + w3 - 3.2, by + 0.1, 3.05, 0.3, anchor=MSO_ANCHOR.TOP)
para(lab, "PROPOSAL DECK · CIO COMMENTARY", 9, FAINT, bold=True, first=True,
     align=PP_ALIGN.RIGHT, ls=0.4)

# Output C — proposed allocation (v7)
cy2 = by + oh + gap
b = box(x3, cy2, w3, oh, fill=GREEN_TINT, line=GREEN, line_w=1.25, radius=0.08)
tf = b.text_frame
para(tf, "\U0001F4E5  Proposed allocation", 13, INK, bold=True, first=True)
para(tf, "If the client stated weights → offered to fill the sleeves.", 10.5, SOFT,
     before=3, spacing=1.04)
para(tf, "Nasdaq 20% + S&P 30% → equity 50%", 10.5, SOFT, font=MONO, before=4)
lab = txt(x3 + w3 - 2.4, cy2 + 0.1, 2.28, 0.3, anchor=MSO_ANCHOR.TOP)
para(lab, "YOU APPLY / IGNORE", 9, GREEN, bold=True, first=True, align=PP_ALIGN.RIGHT, ls=0.5)

# small branching chevron into the gap between confirm and the outputs
chevron(x3 - 0.25, SPINE)

# ---- guardrail band ------------------------------------------------------- #
gy = 6.0
b = box(0.8, gy, 11.73, 1.0, fill=NAVY_DEEP, line=None, radius=0.05)
tf = b.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "THE GUARDRAIL · v8 TIERED", 11, GOLD, bold=True, first=True, ls=1.2)
para(tf, "Each item is tiered 🔒 enforced / 📡 monitored / 📝 advisory. An enforced price trigger "
         "can gate a rebalance buy (→ hold) against a sourced live price; everything else shapes "
         "the narrative & watchlist. A client's condition can gate or flag a trade — never invent one.",
     13, WHITE, font=SERIF, italic=True, before=4, spacing=1.08)

OUT = "Tactical_Workflow_OnePager.pptx"
prs.save(OUT)
print(f"wrote {OUT}")
