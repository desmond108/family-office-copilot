"""tactical_workflow_pptx.py — the "after Confirm" tactical-instructions workflow
as a native 16:9 PowerPoint deck, for dropping into partner / client decks.

Six slides that mirror the shareable artifact (tactical_workflow.html): what the
copilot does the moment the analyst confirms the sorted tactical items — the
structured record, the two-stream split (watchlist / guidance), the three places
the items appear, and the guardrail (guidance never moves a number). House
navy / gold style, matching generate_proposal.py.

Run:  python tactical_workflow_pptx.py   ->  Tactical_Workflow.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (matches the artifact + generate_proposal) ------------------- #
NAVY = RGBColor(0x1E, 0x2A, 0x56)
NAVY_DEEP = RGBColor(0x13, 0x1D, 0x3D)
GOLD = RGBColor(0xB0, 0x87, 0x2A)
GOLD_DK = RGBColor(0x9C, 0x74, 0x22)
GOLD_TINT = RGBColor(0xF7, 0xEF, 0xDC)
GREEN = RGBColor(0x2F, 0x6F, 0x57)
GREEN_TINT = RGBColor(0xE6, 0xF1, 0xEB)
INK = RGBColor(0x16, 0x20, 0x3C)
SOFT = RGBColor(0x47, 0x51, 0x7A)
FAINT = RGBColor(0x83, 0x8C, 0xAD)
LINE = RGBColor(0xDF, 0xE2, 0xEE)
LINE_STRONG = RGBColor(0xC8, 0xCD, 0xDD)
SLATE_TINT = RGBColor(0xEE, 0xF0, 0xF6)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xB9, 0xC2, 0xE0)

SERIF = "Georgia"          # closest ubiquitous serif to the artifact's Palatino stack
SANS = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = In(13.333)
prs.slide_height = In(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _no_autosize(tf):
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    # corner radius
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        # keep it simple — no custom shadow XML; flat design reads clean
    tf = shp.text_frame
    tf.margin_left = In(0.16); tf.margin_right = In(0.16)
    tf.margin_top = In(0.12); tf.margin_bottom = In(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _no_autosize(tf)
    return shp


def txt(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame; _no_autosize(tf)
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, *, font=SANS, bold=False, italic=False,
         first=False, before=0, after=0, align=PP_ALIGN.LEFT, spacing=1.0, ls=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before); p.space_after = Pt(after)
    try:
        p.line_spacing = spacing
    except Exception:
        pass
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    if ls is not None:
        _set_letter_spacing(r, ls)
    return p


def _set_letter_spacing(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def snum(s, n, light=False):
    tf = txt(s, 11.3, 0.34, 1.7, 0.4)
    para(tf, f"{n:02d} / 06", 10.5, (CREAM if light else FAINT), font=MONO,
         first=True, align=PP_ALIGN.RIGHT, ls=0.4)


def eyebrow(s, text, x=0.86, y=0.72):
    tf = txt(s, x, y, 11, 0.4)
    para(tf, text.upper(), 11.5, GOLD, font=SANS, bold=True, first=True, ls=1.6)


def title(s, text, x=0.86, y=1.12, w=11.6, size=26, color=NAVY):
    tf = txt(s, x, y, w, 1.4)
    para(tf, text, size, color, font=SERIF, bold=True, first=True, spacing=1.05)
    return tf


def sub(s, text, x=0.86, y=2.02, w=11.2, size=14):
    tf = txt(s, x, y, w, 0.7)
    para(tf, text, size, SOFT, first=True, spacing=1.1)


# =========================================================================== #
# Slide 1 — title
# =========================================================================== #
s = slide(NAVY_DEEP)
box(s, 0, 0, 0.14, 7.5, fill=GOLD)                    # gold rail
snum(s, 1, light=True)
eyebrow(s, "The tactical-instructions workflow", x=0.9, y=1.4)
tf = txt(s, 0.9, 2.4, 10.6, 2.4)
para(tf, "From the client's words to a working proposal", 38, WHITE,
     font=SERIF, bold=True, first=True, spacing=1.02)
tf = txt(s, 0.9, 4.7, 9.6, 1.0)
para(tf, "What the copilot does the moment you confirm the sorted items — and the "
         "one rule that keeps every number trustworthy.", 15.5, CREAM, first=True, spacing=1.2)
tf = txt(s, 0.9, 6.5, 11.5, 0.5)
para(tf, "Meridian Family Office Copilot   ·   v8 · Tactical instructions   ·   "
         "For partner & client discussion", 11.5, RGBColor(0x7F, 0x8C, 0xBB), first=True)


# =========================================================================== #
# Slide 2 — pipeline at a glance
# =========================================================================== #
s = slide()
snum(s, 2)
eyebrow(s, "At a glance")
title(s, "Four steps in, three outputs out")
sub(s, "The first four steps are what the analyst does. Confirm is where this walkthrough "
       "begins — everything to its right is automatic.")

nodes = [
    ("01", "Paste", "Client's asks, plain language", SLATE_TINT, INK, LINE),
    ("02", "Sort", "Copilot types each ask", SLATE_TINT, INK, LINE),
    ("03", "Review", "Keep · edit · flag unclear", SLATE_TINT, INK, LINE),
    ("04", "Confirm", "You are here", GREEN_TINT, GREEN, GREEN),
    ("→", "Watchlist", "Levels to monitor", GOLD_TINT, GOLD_DK, GOLD),
    ("→", "Guidance", "Shapes the proposal", GOLD_TINT, GOLD_DK, GOLD),
    ("→", "Allocation", "If weights → Apply", GOLD_TINT, GOLD_DK, GOLD),
]
n = len(nodes); gap = 0.22; x0 = 0.86; total = 11.6
w = (total - gap * (n - 1)) / n
y = 3.1; h = 1.9
for i, (k, t, d, fill, tcol, lcol) in enumerate(nodes):
    x = x0 + i * (w + gap)
    b = box(s, x, y, w, h, fill=fill, line=lcol, line_w=1.25, radius=0.10)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP
    para(tf, k, 12, FAINT if fill == SLATE_TINT else tcol, font=MONO, bold=True, first=True, ls=0.4)
    para(tf, ("✓ " + t) if t == "Confirm" else t, 15, tcol, bold=True, before=4, spacing=1.0)
    para(tf, d, 11, SOFT, before=3, spacing=1.05)
    if i < n - 1:
        cx = x + w + (gap - 0.16) / 2
        ar = txt(s, cx - 0.05, y + h / 2 - 0.2, 0.3, 0.4)
        para(ar, "›", 20, LINE_STRONG, first=True, align=PP_ALIGN.CENTER)

tf = txt(s, 0.86, 5.35, 11, 0.5)
para(tf, "The rest of this deck follows what leaves the Confirm step.", 13, SOFT,
     first=True, bold=True)


# =========================================================================== #
# Slide 3 — structured record
# =========================================================================== #
s = slide()
snum(s, 3)
eyebrow(s, "Step 1 · after Confirm")
title(s, "Each ask is now a structured record")
sub(s, "A sentence becomes fields the system can act on. Every level is copied from the "
       "client's own words — never invented.")

# outer gold-tinted record frame
box(s, 0.86, 3.0, 11.6, 1.75, fill=GOLD_TINT, line=GOLD, line_w=1.5, radius=0.06)
fields = [
    ("TYPE", "entry_trigger", GOLD_DK, MONO),
    ("INSTRUMENT", "Gold ETF", INK, MONO),
    ("ACTION", "Buy", INK, MONO),
    ("LEVEL (COPIED)", "USD 4,000/oz", GOLD_DK, MONO),
]
fw = (11.6 - 0.18 - 0.14 * (len(fields) - 1)) / len(fields)
fx = 0.86 + 0.09; fy = 3.16
for i, (k, v, vcol, vfont) in enumerate(fields):
    x = fx + i * (fw + 0.14)
    b = box(s, x, fy, fw, 1.43, fill=PAPER, line=LINE, line_w=1.0, radius=0.09)
    tf = b.text_frame
    para(tf, k, 10, FAINT, bold=True, first=True, ls=0.8)
    para(tf, v, 15, vcol, font=vfont, bold=True, before=8, spacing=1.05)

tf = txt(s, 0.86, 5.05, 11.6, 0.9)
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'From:  '
r.font.size = Pt(12.5); r.font.name = SANS; r.font.color.rgb = SOFT
r = p.add_run(); r.text = '"the gold ETF can be bought below USD 4,000/oz"'
r.font.size = Pt(12.5); r.font.name = MONO; r.font.color.rgb = INK
para(tf, "Nothing computed — just sorted and copied.", 12.5, GREEN, bold=True, before=6)


# =========================================================================== #
# Slide 4 — two streams
# =========================================================================== #
s = slide()
snum(s, 4)
eyebrow(s, "Step 2 · the split")
title(s, "The items split by what they're for")

sw = 5.68; sh = 3.15; sy = 2.75
# Stream A — triggers -> watchlist
b = box(s, 0.86, sy, sw, sh, fill=GOLD_TINT, line=GOLD, line_w=1.5, radius=0.05)
tf = b.text_frame
para(tf, "\U0001F4E1   Entry triggers  →  Monitoring watchlist", 15, INK, bold=True, first=True)
para(tf, "Conditional, level-based asks become a live list the portfolio is watched "
         "against — the retention hook.", 12.8, SOFT, before=8, spacing=1.15)
inner = box(s, 1.12, sy + 1.72, sw - 0.52, 1.1, fill=PAPER, line=LINE_STRONG, line_w=1.0, radius=0.08)
tfi = inner.text_frame
p = tfi.paragraphs[0]
for seg, col, bold in [("Gold ETF · below ", SOFT, False), ("USD 4,000/oz", GOLD_DK, True),
                       (" · Buy", SOFT, False)]:
    r = p.add_run(); r.text = seg; r.font.size = Pt(11.5); r.font.name = MONO
    r.font.color.rgb = col; r.font.bold = bold
p2 = tfi.add_paragraph()
for seg, col, bold in [("S&P 500 ETF · ", SOFT, False), ("−15% to −20%", GOLD_DK, True),
                       (" from high", SOFT, False)]:
    r = p2.add_run(); r.text = seg; r.font.size = Pt(11.5); r.font.name = MONO
    r.font.color.rgb = col; r.font.bold = bold

# Stream B — every item -> guidance
b = box(s, 0.86 + sw + 0.24, sy, sw, sh, fill=SLATE_TINT, line=LINE_STRONG, line_w=1.25, radius=0.05)
tf = b.text_frame
para(tf, "\U0001F9ED   Every item  →  Analyst guidance", 15, INK, bold=True, first=True)
para(tf, "All items (including the triggers) carry forward as intent that shapes the "
         "written advice — never as figures.", 12.8, SOFT, before=8, spacing=1.15)
inner = box(s, 0.86 + sw + 0.24 + 0.26, sy + 1.72, sw - 0.52, 1.1, fill=PAPER,
            line=LINE_STRONG, line_w=1.0, radius=0.08)
tfi = inner.text_frame
for i, line in enumerate(['Execution style · "buy in tranches"',
                          'Selection criteria · "low fees, good liquidity"',
                          'Open question · "impact of rate hikes?"']):
    para(tfi, line, 11, SOFT, font=MONO, first=(i == 0), before=(0 if i == 0 else 3), spacing=1.1)

# note bar — clarification hold-out + allocation proposal
nb = box(s, 0.86, 6.1, 11.6, 0.75, fill=SLATE_TINT, line=LINE, line_w=1.0, radius=0.06)
tf = nb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "⚠️  Anything ambiguous is typed needs-clarification and held out of the "
         "proposal until resolved. If the client stated target weights, Confirm surfaces a "
         "Proposed allocation to Apply.", 12, SOFT, first=True, spacing=1.12)


# =========================================================================== #
# Slide 5 — three destinations
# =========================================================================== #
s = slide()
snum(s, 5)
eyebrow(s, "Step 3 · where they appear")
title(s, "Three places the analyst sees them")

dw = 3.72; dh = 3.55; dy = 2.65; dx0 = 0.86; dgap = 0.22
dests = [
    ("\U0001F4E1", "Monitoring watchlist", "Triggers, ready to watch the book against.",
     "IN THE COPILOT", GOLD_DK, GOLD, True),
    ("\U0001F4C4", "Proposal deck", 'Listed under "Analyst notes folded into this proposal" '
     "on the data & method slide.", "PPTX · PDF", INK, LINE, False),
    ("\U0001F4AC", "CIO commentary", "Prose shaped by the guidance — quoting only the "
     "computed figures.", "GENERATED NARRATIVE", INK, LINE, False),
]
for i, (ic, t, d, where, tcol, lcol, tri) in enumerate(dests):
    x = dx0 + i * (dw + dgap)
    b = box(s, x, dy, dw, dh, fill=PAPER, line=lcol, line_w=1.5 if tri else 1.0, radius=0.06)
    tf = b.text_frame
    para(tf, ic, 22, INK, first=True)
    para(tf, t, 14, tcol, bold=True, before=8)
    para(tf, d, 12, SOFT, before=6, spacing=1.15)
    if tri:
        p = tf.add_paragraph(); p.space_before = Pt(10)
        for seg, col in [("S&P 500 ETF  ", INK), ("−15/−20%", GOLD_DK)]:
            r = p.add_run(); r.text = seg; r.font.size = Pt(11); r.font.name = MONO; r.font.color.rgb = col
            r.font.bold = (col == INK)
        p = tf.add_paragraph()
        for seg, col in [("Gold ETF  ", INK), ("< $4,000", GOLD_DK)]:
            r = p.add_run(); r.text = seg; r.font.size = Pt(11); r.font.name = MONO; r.font.color.rgb = col
            r.font.bold = (col == INK)
    lab = txt(s, x + 0.16, dy + dh - 0.5, dw - 0.32, 0.35)
    para(lab, where, 10, FAINT, bold=True, first=True, ls=0.6)


# =========================================================================== #
# Slide 6 — guardrail
# =========================================================================== #
s = slide()
snum(s, 6)
eyebrow(s, "Step 4 · the guardrail (v8: tiered)")
title(s, "Guidance can gate a number — never invent one")

gw = 5.68; gh = 2.45; gy = 2.5
b = box(s, 0.86, gy, gw, gh, fill=GOLD_TINT, line=GOLD, line_w=1.5, radius=0.06)
tf = b.text_frame
para(tf, "GUIDANCE CAN SHAPE & GATE", 11, GOLD_DK, bold=True, first=True, ls=0.8)
for i, li in enumerate(["The monitoring watchlist & the commentary",
                        "🔒 Enforced trigger gates a rebalance buy (→ hold) vs a sourced live price",
                        "What the analyst is prompted to weigh"]):
    para(tf, "•  " + li, 12.5, SOFT, before=(10 if i == 0 else 6), spacing=1.1)

b = box(s, 0.86 + gw + 0.24, gy, gw, gh, fill=GREEN_TINT, line=GREEN, line_w=1.5, radius=0.06)
tf = b.text_frame
para(tf, "GUIDANCE NEVER INVENTS", 11, GREEN, bold=True, first=True, ls=0.8)
for i, li in enumerate(["A figure from thin air — enforced checks use a price with provenance",
                        "Suitability thresholds",
                        "Holdings — every dollar from the real book"]):
    para(tf, "•  " + li, 12.5, SOFT, before=(10 if i == 0 else 6), spacing=1.1)

tf = txt(s, 0.86, gy + gh + 0.45, 11.6, 1.4)
para(tf, "Each item is tiered 🔒 enforced / 📡 monitored / 📝 advisory. A client's condition can "
         "gate or flag a trade — never fabricate one.", 18, NAVY, font=SERIF, italic=True, bold=True,
     first=True, spacing=1.15)


OUT = "Tactical_Workflow.pptx"
prs.save(OUT)
print(f"wrote {OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
